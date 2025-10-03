#!/usr/bin/env python3
"""
MergeCompositor with PowerPoint Shape Integration

Implements the compositor pattern for feMerge filter execution,
integrating with python-pptx shapes and our clean slate architecture.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import List, Dict, Any, Protocol, Optional, TYPE_CHECKING
import logging

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide


# ===== Adapter interfaces =====

class ShapeAdapter(Protocol):
    """
    Shape adapter interface for compositor operations.
    Implement this for your specific PPTX backend.
    """

    def clone(self) -> "ShapeAdapter":
        """Return a deep clone of the shape."""
        ...

    def set_effect_input(self, buffer_name: str) -> None:
        """Route the shape's filter to use input buffer."""
        ...

    def set_opacity(self, alpha_0_100: float) -> None:
        """Set shape alpha (0-100)."""
        ...

    def add_to_parent(self) -> None:
        """Insert shape into document tree."""
        ...

    def place_above(self, other: "ShapeAdapter") -> None:
        """Z-order this shape above another."""
        ...

    def debug_name(self) -> str:
        """Human-readable label."""
        ...


@dataclass(frozen=True)
class MergeRenderStep:
    """Single step in merge render plan."""
    input: str
    z: int  # 0-based, 0 is back


@dataclass(frozen=True)
class MergeRenderPlan:
    """Complete plan for merge rendering."""
    steps: List[MergeRenderStep]


# ===== PowerPoint Shape Adapter =====

class PowerPointShapeAdapter:
    """
    Adapter for python-pptx shapes to work with compositor.
    """

    def __init__(self, shape: 'BaseShape', slide: 'Slide',
                 shape_factory: Optional[Any] = None):
        """
        Initialize adapter with PowerPoint shape.

        Args:
            shape: python-pptx shape to adapt
            slide: Slide containing the shape
            shape_factory: Factory for creating new shapes
        """
        self.shape = shape
        self.slide = slide
        self.shape_factory = shape_factory
        self._effect_input = "SourceGraphic"
        self.logger = logging.getLogger(__name__)

    def clone(self) -> "PowerPointShapeAdapter":
        """Clone shape using PPTX duplicate methods."""
        try:
            # Method 1: Use shape factory if available
            if self.shape_factory:
                cloned_shape = self.shape_factory.duplicate_shape(self.shape)
                return PowerPointShapeAdapter(cloned_shape, self.slide, self.shape_factory)

            # Method 2: Manual clone via XML manipulation
            from lxml import etree
            import copy

            # Get shape XML element
            shape_element = self.shape.element
            cloned_element = copy.deepcopy(shape_element)

            # Add to slide shapes
            self.slide.shapes._spTree.append(cloned_element)

            # Get the new shape from slide
            cloned_shape = self.slide.shapes[-1]

            return PowerPointShapeAdapter(cloned_shape, self.slide, self.shape_factory)

        except Exception as e:
            self.logger.error(f"Failed to clone shape: {e}")
            # Return self as fallback (no clone)
            return self

    def set_effect_input(self, buffer_name: str) -> None:
        """Apply filter buffer to shape effects."""
        self._effect_input = buffer_name

        # In real implementation, this would modify the shape's effect XML
        # For now, store for later application
        if hasattr(self.shape, '_element'):
            # Add custom property to track buffer
            self.shape._element.set('filter-input', buffer_name)

    def set_opacity(self, alpha_0_100: float) -> None:
        """Set shape transparency."""
        try:
            # Convert to PPTX transparency (0=opaque, 100=transparent)
            transparency = 100 - alpha_0_100

            if hasattr(self.shape, 'fill'):
                self.shape.fill.transparency = transparency / 100.0

            # Also set on line if present
            if hasattr(self.shape, 'line') and self.shape.line:
                self.shape.line.transparency = transparency / 100.0

        except Exception as e:
            self.logger.warning(f"Could not set opacity: {e}")

    def add_to_parent(self) -> None:
        """Shape is already in slide, just ensure it's registered."""
        # python-pptx shapes are auto-added when created
        pass

    def place_above(self, other: "PowerPointShapeAdapter") -> None:
        """Reorder shape in z-order."""
        try:
            # Get shape indices
            shapes = list(self.slide.shapes)
            self_idx = shapes.index(self.shape)
            other_idx = shapes.index(other.shape)

            if self_idx <= other_idx:
                return  # Already above

            # Move via XML manipulation
            self_element = self.shape.element
            other_element = other.shape.element
            parent = self_element.getparent()

            # Remove and re-insert after other
            parent.remove(self_element)
            other_idx = list(parent).index(other_element)
            parent.insert(other_idx + 1, self_element)

        except Exception as e:
            self.logger.warning(f"Could not reorder shapes: {e}")

    def debug_name(self) -> str:
        """Get debug name for shape."""
        shape_type = type(self.shape).__name__
        shape_id = getattr(self.shape, 'shape_id', 'unknown')
        return f"{shape_type}#{shape_id}[{self._effect_input}]"


# ===== Main Compositor =====

class MergeCompositor:
    """
    Execute feMerge render plans by composing shape layers.
    """

    def __init__(self, logger: Optional[logging.Logger] = None):
        """Initialize compositor."""
        self.logger = logger or logging.getLogger(__name__)

    def compose(self,
                base_shape: ShapeAdapter,
                plan: MergeRenderPlan,
                *,
                set_layer_alpha: Optional[float] = None,
                return_layers: bool = True) -> List[ShapeAdapter]:
        """
        Compose layers according to plan.

        Args:
            base_shape: Shape to clone for each input
            plan: Render plan from MergeFilter
            set_layer_alpha: Optional alpha for each layer
            return_layers: Whether to return created layers

        Returns:
            List of created shape adapters
        """
        if not plan.steps:
            self.logger.info("MergeCompositor: no steps")
            return []

        created: List[ShapeAdapter] = []
        last_shape: Optional[ShapeAdapter] = None

        # Create clones for each input
        for step in plan.steps:
            clone = base_shape.clone()
            clone.set_effect_input(step.input)

            if set_layer_alpha is not None:
                clone.set_opacity(set_layer_alpha)

            clone.add_to_parent()
            created.append(clone)

            self.logger.debug(
                f"[feMerge] created layer for input='{step.input}' "
                f"z={step.z} as {clone.debug_name()}"
            )

        # Z-order: 0 is back, last is front
        for idx, shape in enumerate(created):
            if last_shape is not None:
                shape.place_above(last_shape)
                self.logger.debug(
                    f"[feMerge] place {shape.debug_name()} "
                    f"above {last_shape.debug_name()}"
                )
            last_shape = shape

        self.logger.info(
            f"MergeCompositor: stacked {len(created)} layers "
            f"(0=back → {len(created)-1}=front)"
        )

        return created if return_layers else []

    @staticmethod
    def plan_from_metadata(metadata: Dict[str, Any]) -> MergeRenderPlan:
        """
        Build plan from FilterResult metadata.

        Args:
            metadata: Metadata with render_plan or merge_inputs

        Returns:
            MergeRenderPlan for execution
        """
        if "render_plan" in metadata and isinstance(metadata["render_plan"], list):
            steps = [
                MergeRenderStep(input=step["input"], z=int(step["z"]))
                for step in metadata["render_plan"]
            ]
            return MergeRenderPlan(steps=steps)

        # Fallback to merge_inputs
        inputs = metadata.get("merge_inputs") or []
        steps = [
            MergeRenderStep(input=name, z=i)
            for i, name in enumerate(inputs)
        ]
        return MergeRenderPlan(steps=steps)


# ===== Factory functions =====

def create_pptx_adapter(shape: 'BaseShape', slide: 'Slide') -> PowerPointShapeAdapter:
    """Create adapter for python-pptx shape."""
    return PowerPointShapeAdapter(shape, slide)


def create_compositor(logger: Optional[logging.Logger] = None) -> MergeCompositor:
    """Create merge compositor instance."""
    return MergeCompositor(logger)


# ===== Integration helper =====

class FilterCompositorIntegration:
    """
    Helper to integrate compositor with filter pipeline.
    """

    def __init__(self, slide: 'Slide', policy: Optional[Any] = None):
        """
        Initialize integration.

        Args:
            slide: PowerPoint slide for shape operations
            policy: Policy engine for decisions
        """
        self.slide = slide
        self.policy = policy
        self.compositor = MergeCompositor()
        self.logger = logging.getLogger(__name__)

    def apply_merge_filter(self, shape: 'BaseShape',
                          filter_result: Dict[str, Any]) -> List['BaseShape']:
        """
        Apply merge filter result to shape.

        Args:
            shape: Base shape to apply filter to
            filter_result: Result from MergeFilter.apply()

        Returns:
            List of created shapes (if layer_stack strategy)
        """
        metadata = filter_result.get("metadata", {})
        strategy = filter_result.get("strategy", "single_composite")

        if strategy == "layer_stack":
            # Use compositor for layer stacking
            adapter = PowerPointShapeAdapter(shape, self.slide)
            plan = MergeCompositor.plan_from_metadata(metadata)

            # Compose layers
            layers = self.compositor.compose(adapter, plan)

            # Return actual shapes
            return [layer.shape for layer in layers]

        elif strategy == "single_composite":
            # Apply simplified effect to single shape
            drawingml = filter_result.get("drawingml")
            if drawingml:
                self._apply_effect_xml(shape, drawingml)
            return [shape]

        else:  # emf_rasterize
            # Would trigger EMF generation pipeline
            self.logger.info(f"EMF rasterization requested for {shape}")
            return [shape]

    def _apply_effect_xml(self, shape: 'BaseShape', effect_xml: str) -> None:
        """Apply effect XML to shape (simplified)."""
        # In real implementation, would parse and apply XML
        self.logger.debug(f"Applying effect to shape: {effect_xml[:100]}...")