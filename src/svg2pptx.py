#!/usr/bin/env python3
"""
SVG to PowerPoint Converter

Complete pipeline for converting SVG files to PowerPoint presentations
with proper DrawingML vector graphics integration.
"""

import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from .svg2drawingml import SVGToDrawingMLConverter

import tempfile
import zipfile
from lxml import etree as ET


class SVGToPowerPointConverter:
    """Convert SVG files to PowerPoint presentations with vector graphics."""
    
    def __init__(self, slide_width=10, slide_height=7.5):
        """
        Initialize converter.
        
        Args:
            slide_width: Slide width in inches (default: 10")
            slide_height: Slide height in inches (default: 7.5")
        """
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.svg_converter = SVGToDrawingMLConverter()
    
    def convert_file(self, svg_file: str, output_file: str = None) -> str:
        """
        Convert SVG file to PowerPoint presentation.
        
        Args:
            svg_file: Path to SVG file
            output_file: Output PPTX file (default: same name as SVG)
        
        Returns:
            Path to created PPTX file
        """
        if not output_file:
            svg_path = Path(svg_file)
            output_file = svg_path.with_suffix('.pptx')
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(self.slide_width)
        prs.slide_height = Inches(self.slide_height)
        
        # Add slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Convert SVG to DrawingML
        try:
            drawingml_shapes = self.svg_converter.convert_file(svg_file)
            
            # Add shapes to slide using low-level XML manipulation
            self._add_drawingml_to_slide(slide, drawingml_shapes)
            
        except Exception as e:
            print(f"Warning: Could not convert SVG shapes: {e}")
            # Fallback: Add a text box with error message
            textbox = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(1)
            )
            textbox.text = f"SVG conversion failed: {svg_file}"
        
        # Save presentation
        prs.save(output_file)
        return str(output_file)
    
    def _add_drawingml_to_slide(self, slide, drawingml_shapes: str):
        """Add DrawingML shapes to slide using XML manipulation."""

        if not drawingml_shapes or not drawingml_shapes.strip():
            return

        drawingml_shapes = drawingml_shapes.strip()

        sp_tree = slide.shapes._spTree
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

        # Collect existing shape IDs and names to avoid collisions on insert.
        existing_ids = set()
        existing_names = set()
        for cNvPr in sp_tree.xpath(".//p:cNvPr", namespaces=nsmap):
            id_attr = cNvPr.get("id")
            try:
                existing_ids.add(int(id_attr))
            except (TypeError, ValueError):
                pass
            name_attr = cNvPr.get("name")
            if name_attr:
                existing_names.add(name_attr)

        next_id = (max(existing_ids) + 1) if existing_ids else 1

        # Determine which namespace prefixes are used so that we can declare them
        # on a temporary root wrapper element for parsing.
        prefix_pattern = re.compile(r"([A-Za-z_][A-Za-z0-9_\-\.]*):[A-Za-z_]")
        detected_prefixes = set(prefix_pattern.findall(drawingml_shapes))

        namespace_map = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
            "a15": "http://schemas.microsoft.com/office/drawing/2012/main",
            "a16": "http://schemas.microsoft.com/office/drawing/2014/main",
            "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
            "p15": "http://schemas.microsoft.com/office/powerpoint/2012/main",
            "p16": "http://schemas.microsoft.com/office/powerpoint/2016/main",
            "p17": "http://schemas.microsoft.com/office/powerpoint/2019/main",
            "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
            "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
            "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
            "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
            "lc": "http://schemas.openxmlformats.org/drawingml/2006/lockedCanvas",
            "o": "urn:schemas-microsoft-com:office:office",
            "v": "urn:schemas-microsoft-com:vml",
            "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
            "wps": "http://schemas.microsoft.com/office/word/2010/wordprocessingShape",
            "wpg": "http://schemas.microsoft.com/office/word/2010/wordprocessingGroup",
            "wpi": "http://schemas.microsoft.com/office/word/2010/wordprocessingInk",
            "wne": "http://schemas.microsoft.com/office/word/2006/wordml",
            "xdr": "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
        }

        # Always include the core namespaces expected by DrawingML.
        detected_prefixes.update({"a", "p", "r"})
        namespace_attrs = " ".join(
            f"xmlns:{prefix}=\"{namespace_map[prefix]}\""
            for prefix in sorted(detected_prefixes)
            if prefix in namespace_map
        )

        wrapper_xml = f"<root {namespace_attrs}>{drawingml_shapes}</root>"
        try:
            wrapper = ET.fromstring(wrapper_xml)
        except ET.XMLSyntaxError as exc:
            raise ValueError("Invalid DrawingML content provided for insertion") from exc

        ext_lst = sp_tree.find(qn('p:extLst'))
        insert_index = len(sp_tree)
        if ext_lst is not None:
            insert_index = list(sp_tree).index(ext_lst)

        for element in list(wrapper):
            if not isinstance(element.tag, str):
                continue

            for cNvPr in element.xpath(".//p:cNvPr", namespaces=nsmap):
                id_attr = cNvPr.get("id")
                assign_new_id = False
                try:
                    id_value = int(id_attr)
                except (TypeError, ValueError):
                    assign_new_id = True
                    id_value = None
                else:
                    if id_value in existing_ids:
                        assign_new_id = True

                if assign_new_id:
                    id_value = next_id
                    next_id += 1
                    cNvPr.set("id", str(id_value))
                else:
                    if id_value >= next_id:
                        next_id = id_value + 1

                existing_ids.add(id_value)

                name_attr = cNvPr.get("name")
                if name_attr and name_attr not in existing_names:
                    existing_names.add(name_attr)
                else:
                    base_name = name_attr or "SVG Shape"
                    suffix = 2 if name_attr else 1
                    candidate = base_name if not name_attr else f"{base_name} {suffix}"
                    while candidate in existing_names:
                        suffix += 1
                        candidate = f"{base_name} {suffix}"
                    cNvPr.set("name", candidate)
                    existing_names.add(candidate)

            if ext_lst is not None:
                sp_tree.insert(insert_index, element)
                insert_index += 1
            else:
                sp_tree.append(element)
    
    def batch_convert(self, svg_directory: str, output_directory: str = None):
        """
        Convert all SVG files in a directory.
        
        Args:
            svg_directory: Directory containing SVG files
            output_directory: Output directory (default: same as input)
        """
        svg_dir = Path(svg_directory)
        output_dir = Path(output_directory) if output_directory else svg_dir
        
        svg_files = list(svg_dir.glob('*.svg'))
        if not svg_files:
            print(f"No SVG files found in {svg_directory}")
            return
        
        output_dir.mkdir(exist_ok=True)
        
        for svg_file in svg_files:
            output_file = output_dir / f"{svg_file.stem}.pptx"
            print(f"Converting {svg_file.name} -> {output_file.name}")
            
            try:
                self.convert_file(str(svg_file), str(output_file))
                print(f"  ✓ Created {output_file}")
            except Exception as e:
                print(f"  ✗ Error: {e}")


def main():
    """Command-line interface."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Convert SVG files to PowerPoint presentations')
    parser.add_argument('input', help='SVG file or directory to convert')
    parser.add_argument('-o', '--output', help='Output file or directory')
    parser.add_argument('--width', type=float, default=10, help='Slide width in inches (default: 10)')
    parser.add_argument('--height', type=float, default=7.5, help='Slide height in inches (default: 7.5)')
    parser.add_argument('--batch', action='store_true', help='Batch convert all SVG files in directory')
    
    args = parser.parse_args()
    
    # Create converter
    converter = SVGToPowerPointConverter(args.width, args.height)
    
    if args.batch or Path(args.input).is_dir():
        # Batch conversion
        converter.batch_convert(args.input, args.output)
    else:
        # Single file conversion
        if not Path(args.input).exists():
            print(f"Error: File not found: {args.input}")
            sys.exit(1)
        
        output_file = converter.convert_file(args.input, args.output)
        print(f"✓ Created PowerPoint: {output_file}")


# API Functions for integration tests and external use
def convert_svg_to_pptx(svg_input, output_path: str = None, 
                       slide_width: float = 10, slide_height: float = 7.5,
                       preprocessing_config: dict = None,
                       title: str = None, author: str = None) -> str:
    """
    Convert SVG content or file to PowerPoint presentation.
    
    Args:
        svg_input: SVG content as string OR path to SVG file
        output_path: Output PPTX file path (optional, creates temp file if None)
        slide_width: Slide width in inches (default: 10")
        slide_height: Slide height in inches (default: 7.5")
        preprocessing_config: Optional preprocessing configuration
    
    Returns:
        Path to created PPTX file
    """
    import tempfile
    
    # Determine if input is file path or content
    if isinstance(svg_input, str) and svg_input.endswith('.svg') and os.path.exists(svg_input):
        # Input is a file path
        with open(svg_input, 'r', encoding='utf-8') as f:
            svg_content = f.read()
        temp_svg_path = svg_input  # Use existing file
        cleanup_temp_svg = False
    else:
        # Input is SVG content
        svg_content = svg_input
        # Create temporary SVG file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.svg', delete=False, encoding='utf-8') as f:
            f.write(svg_content)
            temp_svg_path = f.name
        cleanup_temp_svg = True
    
    try:
        # Apply preprocessing if configured
        if preprocessing_config:
            from src.preprocessing import create_optimizer
            optimizer = create_optimizer(**preprocessing_config)
            optimized_content = optimizer.optimize(svg_content)
            
            # Write optimized content back to temp file
            with open(temp_svg_path, 'w', encoding='utf-8') as f:
                f.write(optimized_content)
        
        # Create converter and convert
        converter = SVGToPowerPointConverter(slide_width, slide_height)
        
        # Generate output path if not provided
        if not output_path:
            output_path = tempfile.mktemp(suffix='.pptx')
        
        result_path = converter.convert_file(temp_svg_path, output_path)
        return result_path
        
    finally:
        # Clean up temporary SVG file only if we created it
        if cleanup_temp_svg:
            try:
                os.unlink(temp_svg_path)
            except OSError:
                pass


if __name__ == "__main__":
    main()