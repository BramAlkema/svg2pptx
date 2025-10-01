#!/usr/bin/env python3
"""
SVG to PowerPoint Converter

Complete pipeline for converting SVG files to PowerPoint presentations
with proper DrawingML vector graphics integration.

CRITICAL ISSUES STATUS
======================
✅ Issue 1: Fix ConversionContext Constructor - COMPLETED
✅ Issue 2: Fix CLI Integration DrawingML - COMPLETED
✅ Issue 3: Fix PPTXBuilder Placement - COMPLETED
✅ Issue 4: Fix NumPy Converter Adapters - COMPLETED
✅ Issue 5: Fix Gradient & Pattern Fill - COMPLETED

🔄 Issue 6: Fix Filter Processing Stack - IN PROGRESS
   See: src/converters/filters/converter.py for details

⏳ Issue 7: Fix Image Conversion - PENDING
   See: src/converters/image.py for details

⏳ Issue 8: Fix Font Embedding - PENDING
   See: src/converters/font_embedding.py for details

⏳ Issue 9: Fix Integration Tests - PENDING
   See: tests/integration/pipeline/test_real_conversion_pipeline.py for details

✅ Issue 10: Multi-Slide Conversion - COMPLETED
   Replaced with Clean Slate multipage system in core/multipage/

CURRENT STATUS: 6/10 critical issues resolved, core functionality working
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from .svg2drawingml import SVGToDrawingMLConverter
from .services.conversion_services import ConversionServices
from .core.pptx_builder import PPTXBuilder
# Old multislide import removed - replaced with Clean Slate multipage system
import tempfile
import zipfile
from lxml import etree as ET

# Import CLI visual reporting (optional - graceful degradation if not available)
try:
    from .cli.visual_reports import CLIVisualReportCoordinator
    CLI_VISUAL_REPORTS_AVAILABLE = True
except ImportError:
    CLI_VISUAL_REPORTS_AVAILABLE = False


class SVGToPowerPointConverter:
    """Convert SVG files to PowerPoint presentations with vector graphics."""
    
    def __init__(self, slide_width=10, slide_height=7.5, services=None):
        """
        Initialize converter.

        Args:
            slide_width: Slide width in inches (default: 10")
            slide_height: Slide height in inches (default: 7.5")
            services: ConversionServices instance (optional, creates default if None)
        """
        self.slide_width = slide_width
        self.slide_height = slide_height

        # Use provided services or create default
        if services is None:
            services = ConversionServices.create_default()
        self.services = services

        self.svg_converter = SVGToDrawingMLConverter(services=services)
        self.pptx_builder = PPTXBuilder()
    
    def convert_file(self, svg_file: str, output_file: str = None) -> str:
        """
        Convert SVG file to PowerPoint presentation.

        Args:
            svg_file: Path to SVG file
            output_file: Output PPTX file (default: same name as SVG)

        Returns:
            Path to created PPTX file
        """
        if not output_file:
            svg_path = Path(svg_file)
            output_file = svg_path.with_suffix('.pptx')

        try:
            # Convert SVG to DrawingML
            drawingml_shapes = self.svg_converter.convert_file(svg_file)

            # Create PPTX file with DrawingML shapes using PPTXBuilder
            self.pptx_builder.create_minimal_pptx(drawingml_shapes, str(output_file))

            return str(output_file)

        except Exception as e:
            print(f"Error: Could not convert SVG file: {e}")
            # Fallback: Create PPTX with python-pptx and error message
            prs = Presentation()
            prs.slide_width = Inches(self.slide_width)
            prs.slide_height = Inches(self.slide_height)

            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            textbox = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(1)
            )
            textbox.text = f"SVG conversion failed: {svg_file}\nError: {e}"

            prs.save(output_file)
            return str(output_file)
    
    
    def batch_convert(self, svg_directory: str, output_directory: str = None):
        """
        Convert all SVG files in a directory.
        
        Args:
            svg_directory: Directory containing SVG files
            output_directory: Output directory (default: same as input)
        """
        svg_dir = Path(svg_directory)
        output_dir = Path(output_directory) if output_directory else svg_dir
        
        svg_files = list(svg_dir.glob('*.svg'))
        if not svg_files:
            print(f"No SVG files found in {svg_directory}")
            return
        
        output_dir.mkdir(exist_ok=True)
        
        for svg_file in svg_files:
            output_file = output_dir / f"{svg_file.stem}.pptx"
            print(f"Converting {svg_file.name} -> {output_file.name}")
            
            try:
                self.convert_file(str(svg_file), str(output_file))
                print(f"  ✓ Created {output_file}")
            except Exception as e:
                print(f"  ✗ Error: {e}")


def _generate_visual_report_if_requested(args, svg_file_path: str, pptx_file_path: str) -> bool:
    """
    Generate visual report if CLI flags request it.

    Args:
        args: Parsed CLI arguments
        svg_file_path: Path to source SVG file
        pptx_file_path: Path to generated PPTX file

    Returns:
        True if visual report generation succeeded or was not requested, False if failed
    """
    # Check if any visual report flags are set
    if not (args.visual_report or args.debug_report or args.open_report or args.google_slides_link):
        return True  # No visual report requested - success

    # Check if visual reporting is available
    if not CLI_VISUAL_REPORTS_AVAILABLE:
        print("⚠️  Visual reporting not available - CLI visual reports module not found")
        return True  # Graceful degradation - don't fail conversion

    try:
        print("📊 Generating visual report...")

        # Initialize visual report coordinator
        coordinator = CLIVisualReportCoordinator(
            output_dir=Path(args.report_output_dir),
            enable_google_slides=args.google_slides_link
        )

        # Generate visual report
        report_path = coordinator.generate_visual_report(
            svg_file=Path(svg_file_path),
            pptx_file=Path(pptx_file_path),
            include_debug=args.debug_report,
            auto_open=args.open_report
        )

        if report_path:
            print(f"✅ Visual report generated: {report_path}")
            return True
        else:
            print("⚠️  Visual report generation failed, but conversion completed successfully")
            return True  # Don't fail conversion due to visual report issues

    except Exception as e:
        print(f"⚠️  Visual report generation error: {e}")
        print("   Conversion completed successfully, but visual report failed")
        return True  # Graceful degradation - don't fail conversion


def main():
    """Command-line interface."""
    import argparse

    parser = argparse.ArgumentParser(description='Convert SVG files to PowerPoint presentations')
    parser.add_argument('input', help='SVG file or directory to convert')
    parser.add_argument('-o', '--output', help='Output file or directory')
    parser.add_argument('--width', type=float, default=10, help='Slide width in inches (default: 10)')
    parser.add_argument('--height', type=float, default=7.5, help='Slide height in inches (default: 7.5)')
    parser.add_argument('--batch', action='store_true', help='Batch convert all SVG files in directory')
    # Multislide options removed - use Clean Slate multipage system instead

    # Visual reporting arguments
    visual_group = parser.add_argument_group('visual reporting', 'Generate visual comparison reports')
    visual_group.add_argument('--visual-report', action='store_true',
                             help='Generate side-by-side visual comparison report with browser and PPTX screenshots')
    visual_group.add_argument('--debug-report', action='store_true',
                             help='Include comprehensive debug logging and performance metrics in visual report')
    visual_group.add_argument('--open-report', action='store_true',
                             help='Automatically open generated visual report in default browser')
    visual_group.add_argument('--google-slides-link', action='store_true',
                             help='Include "Open in Google Slides" button in visual report (requires API credentials)')
    visual_group.add_argument('--report-output-dir', type=str, default='./reports',
                             help='Custom directory for visual report outputs (default: ./reports/)')

    args = parser.parse_args()

    # Validate visual report arguments
    if args.visual_report or args.debug_report or args.open_report or args.google_slides_link:
        # Validate report output directory
        try:
            report_dir = Path(args.report_output_dir)
            if not report_dir.exists():
                report_dir.mkdir(parents=True, exist_ok=True)
                print(f"Created report directory: {report_dir}")
            elif not report_dir.is_dir():
                print(f"Error: Report output path exists but is not a directory: {report_dir}")
                sys.exit(1)
        except PermissionError:
            print(f"Error: Permission denied creating report directory: {args.report_output_dir}")
            sys.exit(1)
        except Exception as e:
            print(f"Error: Cannot create report directory {args.report_output_dir}: {e}")
            sys.exit(1)

    # Handle multi-slide conversion for multiple files
    # Old multislide handling removed - use Clean Slate multipage system instead

    # Create converter with services
    services = ConversionServices.create_default()

    if args.batch or Path(args.input).is_dir():
        # Batch conversion
        converter = SVGToPowerPointConverter(args.width, args.height, services=services)
        converter.batch_convert(args.input, args.output)

        # Note: Batch visual reports will be implemented in Task 4.3
        if args.visual_report or args.debug_report or args.open_report or args.google_slides_link:
            print("ℹ️  Batch visual reports not yet implemented (Task 4.3)")
            print("   Use single file conversion with --visual-report for visual reports")
    else:
        # Traditional single file conversion
        converter = SVGToPowerPointConverter(args.width, args.height, services=services)
        if not Path(args.input).exists():
            print(f"Error: File not found: {args.input}")
            sys.exit(1)

        output_file = converter.convert_file(args.input, args.output)
        print(f"✓ Created PowerPoint: {output_file}")

        # Generate visual report if requested
        _generate_visual_report_if_requested(args, args.input, output_file)


# API Functions for integration tests and external use
def convert_svg_to_pptx(svg_input, output_path: str = None,
                       slide_width: float = 10, slide_height: float = 7.5,
                       preprocessing_config: dict = None,
                       title: str = None, author: str = None) -> str:
    """
    Convert SVG content or file to PowerPoint presentation with real DrawingML integration.

    Args:
        svg_input: SVG content as string OR path to SVG file
        output_path: Output PPTX file path (optional, creates temp file if None)
        slide_width: Slide width in inches (default: 10")
        slide_height: Slide height in inches (default: 7.5")
        preprocessing_config: Optional preprocessing configuration
        title: Presentation title (deprecated, ignored)
        author: Presentation author (deprecated, ignored)

    Returns:
        Path to created PPTX file
    """
    import tempfile

    # Determine if input is file path or content
    if isinstance(svg_input, str) and svg_input.lower().endswith('.svg') and os.path.exists(svg_input):
        # Input is a file path
        with open(svg_input, 'r', encoding='utf-8') as f:
            svg_content = f.read()
        temp_svg_path = svg_input  # Use existing file
        cleanup_temp_svg = False
    else:
        # Input is SVG content
        svg_content = svg_input
        # Create temporary SVG file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.svg', delete=False, encoding='utf-8') as f:
            f.write(svg_content)
            temp_svg_path = f.name
        cleanup_temp_svg = True

    try:
        # Apply preprocessing if configured
        if preprocessing_config:
            from .preprocessing import create_optimizer
            optimizer = create_optimizer(**preprocessing_config)
            optimized_content = optimizer.optimize(svg_content)

            # FIXED: Always create separate temp file for preprocessing output
            # =================================================================
            # Create a new temporary file for preprocessed content to avoid
            # overwriting the original user file when input is a file path.
            with tempfile.NamedTemporaryFile(mode='w', suffix='.svg', delete=False, encoding='utf-8') as f:
                f.write(optimized_content)
                # Update temp_svg_path to point to preprocessed content
                # and ensure it gets cleaned up
                if not cleanup_temp_svg:
                    # Original was a file path, now we have a temp file to clean up
                    cleanup_temp_svg = True
                temp_svg_path = f.name

        # Create converter with services and convert using PPTXBuilder
        services = ConversionServices.create_default()
        converter = SVGToPowerPointConverter(slide_width, slide_height, services=services)

        # Generate output path if not provided using secure temporary file
        if not output_path:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                output_path = tmp.name

        result_path = converter.convert_file(temp_svg_path, output_path)
        return result_path

    finally:
        # Clean up temporary SVG file only if we created it
        if cleanup_temp_svg:
            try:
                os.unlink(temp_svg_path)
            except OSError:
                pass




if __name__ == "__main__":
    main()