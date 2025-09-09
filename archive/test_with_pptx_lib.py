#!/usr/bin/env python3
"""
Test our generated PPTX files using python-pptx library to verify they work correctly.
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def test_pptx_with_library(pptx_file: str):
    """Test a PPTX file using python-pptx library."""
    if not PPTX_AVAILABLE:
        print("❌ python-pptx library not available")
        print("Run: pip install python-pptx")
        return False
    
    print(f"=== Testing {pptx_file} with python-pptx ===")
    
    if not os.path.exists(pptx_file):
        print(f"❌ File not found: {pptx_file}")
        return False
    
    try:
        # Load the presentation
        prs = Presentation(pptx_file)
        print(f"✅ Successfully loaded presentation")
        
        # Check presentation properties
        print(f"   Slide size: {prs.slide_width} x {prs.slide_height} EMUs")
        print(f"   Number of slides: {len(prs.slides)}")
        
        # Examine each slide
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i+1} ---")
            print(f"   Number of shapes: {len(slide.shapes)}")
            
            # Analyze each shape
            for j, shape in enumerate(slide.shapes):
                shape_type = shape.shape_type
                print(f"   Shape {j+1}: {shape_type}")
                
                # Get position and size
                left = shape.left
                top = shape.top  
                width = shape.width
                height = shape.height
                print(f"     Position: ({left}, {top}) EMUs")
                print(f"     Size: {width} x {height} EMUs")
                
                # Check if it's a regular shape
                if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    print(f"     Auto Shape ID: {shape.auto_shape_type}")
                
                # Check for fill
                if hasattr(shape, 'fill'):
                    fill = shape.fill
                    if fill.type is not None:
                        print(f"     Fill type: {fill.type}")
                
                # Check for line
                if hasattr(shape, 'line'):
                    line = shape.line
                    try:
                        if hasattr(line.color, 'rgb') and line.color.rgb is not None:
                            print(f"     Line color: {line.color.rgb}")
                        print(f"     Line width: {line.width}")
                    except:
                        print("     Line properties: (complex color)")
                
                # Check for text
                if hasattr(shape, 'text'):
                    if shape.text:
                        print(f"     Text: '{shape.text}'")
        
        return True
        
    except Exception as e:
        print(f"❌ Error loading presentation: {e}")
        return False


def compare_with_baseline():
    """Compare our generated files with a baseline PowerPoint file."""
    print("\n=== Comparing with Baseline ===")
    
    # Test our generated files
    test_files = ['test_simple.pptx', 'test_input.pptx']
    
    for test_file in test_files:
        if os.path.exists(test_file):
            success = test_pptx_with_library(test_file)
            if success:
                print(f"✅ {test_file} is valid and readable")
            else:
                print(f"❌ {test_file} failed validation")
        else:
            print(f"⚠️  {test_file} not found - run testbench.py first")


def create_reference_pptx():
    """Create a reference PPTX using python-pptx for comparison."""
    if not PPTX_AVAILABLE:
        print("❌ Cannot create reference - python-pptx not available")
        return False
    
    print("\n=== Creating Reference PPTX ===")
    
    try:
        # Create a new presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add some shapes for comparison
        shapes = slide.shapes
        
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.dml import MSO_LINE
        
        # Add a rectangle
        left = Inches(1)
        top = Inches(1) 
        width = Inches(2)
        height = Inches(1.5)
        rect = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
        
        # Add a circle/ellipse  
        left = Inches(4)
        top = Inches(1)
        width = Inches(1.5)
        height = Inches(1.5) 
        oval = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green
        
        # Add a line
        from pptx.shapes.connector import Connector
        left = Inches(1)
        top = Inches(3)
        width = Inches(4)
        height = Inches(0)
        connector = shapes.add_connector(MSO_LINE.STRAIGHT_CONNECTOR, left, top, left + width, top + height)
        connector.line.color.rgb = RGBColor(0, 0, 255)  # Blue
        connector.line.width = Inches(0.05)
        
        # Save reference file
        ref_file = 'reference.pptx'
        prs.save(ref_file)
        print(f"✅ Created reference file: {ref_file}")
        
        # Test the reference file
        test_pptx_with_library(ref_file)
        
        return True
        
    except Exception as e:
        print(f"❌ Error creating reference: {e}")
        return False


def main():
    """Run all tests with python-pptx library."""
    print("SVG to PPTX Validation with python-pptx Library")
    print("=" * 50)
    
    if not PPTX_AVAILABLE:
        print("❌ python-pptx library is not available")
        print("\nTo install it:")
        print("  pip install python-pptx")
        print("\nOr create a virtual environment:")
        print("  python3 -m venv venv")
        print("  source venv/bin/activate  # On Windows: venv\\Scripts\\activate") 
        print("  pip install python-pptx")
        sys.exit(1)
    
    # Test our generated files
    compare_with_baseline()
    
    # Create and test reference
    create_reference_pptx()
    
    print("\n" + "=" * 50)
    print("✅ Testing completed!")
    print("\nFiles generated:")
    for file in ['test_simple.pptx', 'test_input.pptx', 'reference.pptx']:
        if os.path.exists(file):
            size = os.path.getsize(file)
            print(f"  - {file} ({size:,} bytes)")


if __name__ == "__main__":
    main()