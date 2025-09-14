"""
File and path related fixtures for testing.

Provides fixtures for creating and managing test files.
"""
from pathlib import Path
from typing import List, Generator
import json
import zipfile

import pytest


@pytest.fixture
def test_data_dir() -> Path:
    """Get path to test data directory.
    
    Returns:
        Path to test data directory.
    """
    return Path(__file__).parent.parent / "data"


@pytest.fixture
def svg_test_files(test_data_dir: Path) -> List[Path]:
    """Get list of SVG test files.
    
    Args:
        test_data_dir: Test data directory path
        
    Returns:
        List of paths to SVG test files.
    """
    svg_dir = test_data_dir / "svg"
    if svg_dir.exists():
        return list(svg_dir.glob("*.svg"))
    return []


@pytest.fixture
def create_test_zip(temp_dir: Path) -> Generator:
    """Factory fixture for creating test ZIP files.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Yields:
        Function to create ZIP files with test content.
    """
    def _create_zip(name: str, files: dict) -> Path:
        """Create a ZIP file with specified content.
        
        Args:
            name: Name of the ZIP file
            files: Dictionary mapping filenames to content
            
        Returns:
            Path to created ZIP file.
        """
        zip_path = temp_dir / name
        with zipfile.ZipFile(zip_path, 'w') as zf:
            for filename, content in files.items():
                zf.writestr(filename, content)
        return zip_path
    
    yield _create_zip


@pytest.fixture
def create_test_json(temp_dir: Path) -> Generator:
    """Factory fixture for creating test JSON files.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Yields:
        Function to create JSON files with test data.
    """
    def _create_json(name: str, data: dict) -> Path:
        """Create a JSON file with specified data.
        
        Args:
            name: Name of the JSON file
            data: Dictionary to serialize to JSON
            
        Returns:
            Path to created JSON file.
        """
        json_path = temp_dir / name
        with open(json_path, 'w') as f:
            json.dump(data, f, indent=2)
        return json_path
    
    yield _create_json


@pytest.fixture
def create_pptx_file(temp_dir: Path) -> Path:
    """Create a sample PPTX file for testing.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created PPTX file.
    """
    from pptx import Presentation
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test Presentation"
    
    pptx_path = temp_dir / "test.pptx"
    prs.save(pptx_path)
    
    return pptx_path


@pytest.fixture
def expected_output_dir(test_data_dir: Path) -> Path:
    """Get path to expected output directory.
    
    Args:
        test_data_dir: Test data directory path
        
    Returns:
        Path to expected output directory.
    """
    expected_dir = test_data_dir / "expected"
    expected_dir.mkdir(exist_ok=True)
    return expected_dir


@pytest.fixture
def baseline_dir(test_data_dir: Path) -> Path:
    """Get path to baseline comparison directory.
    
    Args:
        test_data_dir: Test data directory path
        
    Returns:
        Path to baseline directory for visual comparisons.
    """
    baseline_dir = test_data_dir / "baselines"
    baseline_dir.mkdir(exist_ok=True)
    return baseline_dir


@pytest.fixture
def batch_input_files(temp_dir: Path, sample_svg_content: str) -> List[Path]:
    """Create multiple input files for batch testing.
    
    Args:
        temp_dir: Temporary directory fixture
        sample_svg_content: Sample SVG content fixture
        
    Returns:
        List of paths to created input files.
    """
    files = []
    for i in range(5):
        file_path = temp_dir / f"input_{i}.svg"
        file_path.write_text(sample_svg_content)
        files.append(file_path)
    
    return files


@pytest.fixture
def config_file(temp_dir: Path) -> Path:
    """Create a test configuration file.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created configuration file.
    """
    config = {
        "conversion": {
            "preserve_aspect_ratio": True,
            "embed_fonts": False,
            "optimize_output": True
        },
        "performance": {
            "cache_enabled": True,
            "max_workers": 4
        }
    }
    
    config_path = temp_dir / "test_config.json"
    with open(config_path, 'w') as f:
        json.dump(config, f, indent=2)
    
    return config_path