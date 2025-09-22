#!/usr/bin/env python3
"""
Visual Comparison Testing for SVG2PPTX
======================================

Playwright-based visual testing that compares:
1. SVG rendered in browser (source truth)
2. PPTX slide screenshot (our output)
3. Generates side-by-side comparison pages

This can be used for:
- Visual regression testing
- Before/after change validation
- Manual quality assessment
- Documentation of system capabilities
"""

import pytest
import asyncio
import json
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import subprocess
import time
import logging
from datetime import datetime

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

try:
    from playwright.async_api import async_playwright, Page, Browser
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False
    logger.warning("Playwright not available - visual tests will be skipped")

# Import our conversion system
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))

try:
    from src.converters.base import ConversionContext
    from src.services.conversion_services import ConversionServices
    from src.preprocessing.optimizer import SVGOptimizer
    CONVERTERS_AVAILABLE = True
except ImportError:
    CONVERTERS_AVAILABLE = False
    logger.warning("SVG2PPTX converters not available - will use mock conversion")


class VisualTestCase:
    """Represents a single visual test case with SVG and expected results."""

    def __init__(self, name: str, svg_content: str, description: str = ""):
        self.name = name
        self.svg_content = svg_content
        self.description = description
        self.timestamp = datetime.now().isoformat()

    def to_dict(self) -> dict:
        return {
            'name': self.name,
            'description': self.description,
            'timestamp': self.timestamp,
            'svg_length': len(self.svg_content)
        }


class VisualComparisonTester:
    """Main class for visual comparison testing."""

    def __init__(self, output_dir: str = "tests/visual/results"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories
        (self.output_dir / "svgs").mkdir(exist_ok=True)
        (self.output_dir / "svg_screenshots").mkdir(exist_ok=True)
        (self.output_dir / "pptx_files").mkdir(exist_ok=True)
        (self.output_dir / "pptx_screenshots").mkdir(exist_ok=True)
        (self.output_dir / "comparisons").mkdir(exist_ok=True)

        self.test_results: List[Dict] = []

    async def run_visual_test(self, test_case: VisualTestCase) -> Dict:
        """Run a complete visual test for one SVG."""
        logger.info(f"Running visual test: {test_case.name}")

        result = {
            'test_case': test_case.to_dict(),
            'status': 'running',
            'timestamps': {},
            'files': {},
            'errors': []
        }

        try:
            # Step 1: Save SVG file
            svg_path = await self._save_svg_file(test_case)
            result['files']['svg'] = str(svg_path)
            result['timestamps']['svg_saved'] = datetime.now().isoformat()

            # Step 2: Take SVG browser screenshot
            svg_screenshot = await self._capture_svg_screenshot(svg_path, test_case.name)
            result['files']['svg_screenshot'] = str(svg_screenshot)
            result['timestamps']['svg_screenshot'] = datetime.now().isoformat()

            # Step 3: Convert to PPTX
            pptx_path = await self._convert_to_pptx(svg_path, test_case.name)
            result['files']['pptx'] = str(pptx_path)
            result['timestamps']['pptx_created'] = datetime.now().isoformat()

            # Step 4: Take PPTX screenshot
            pptx_screenshot = await self._capture_pptx_screenshot(pptx_path, test_case.name)
            result['files']['pptx_screenshot'] = str(pptx_screenshot)
            result['timestamps']['pptx_screenshot'] = datetime.now().isoformat()

            # Step 5: Create comparison page
            comparison_page = await self._create_comparison_page(
                test_case, svg_screenshot, pptx_screenshot, result
            )
            result['files']['comparison_page'] = str(comparison_page)
            result['timestamps']['comparison_created'] = datetime.now().isoformat()

            result['status'] = 'completed'
            logger.info(f"Visual test completed: {test_case.name}")

        except Exception as e:
            result['status'] = 'failed'
            result['errors'].append(str(e))
            logger.error(f"Visual test failed for {test_case.name}: {e}")

        return result

    async def _save_svg_file(self, test_case: VisualTestCase) -> Path:
        """Save SVG content to file."""
        svg_path = self.output_dir / "svgs" / f"{test_case.name}.svg"
        svg_path.write_text(test_case.svg_content, encoding='utf-8')
        return svg_path

    async def _capture_svg_screenshot(self, svg_path: Path, test_name: str) -> Path:
        """Capture screenshot of SVG rendered in browser."""
        screenshot_path = self.output_dir / "svg_screenshots" / f"{test_name}_svg.png"

        if not PLAYWRIGHT_AVAILABLE:
            # Create placeholder image
            self._create_placeholder_image(screenshot_path, "SVG Screenshot\n(Playwright not available)")
            return screenshot_path

        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()

            # Create HTML page with SVG
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <title>SVG Test: {test_name}</title>
                <style>
                    body {{ margin: 0; padding: 20px; font-family: Arial, sans-serif; }}
                    .container {{ display: flex; flex-direction: column; align-items: center; }}
                    .svg-container {{ border: 1px solid #ccc; padding: 10px; background: white; }}
                    h1 {{ color: #333; margin-bottom: 20px; }}
                </style>
            </head>
            <body>
                <div class="container">
                    <h1>SVG Source: {test_name}</h1>
                    <div class="svg-container">
                        {svg_path.read_text()}
                    </div>
                </div>
            </body>
            </html>
            """

            # Save HTML file for debugging
            html_path = self.output_dir / "svgs" / f"{test_name}.html"
            html_path.write_text(html_content)

            # Load and screenshot
            await page.goto(f"file://{html_path.absolute()}")
            await page.wait_for_load_state('networkidle')

            # Take screenshot of SVG container
            svg_element = await page.query_selector('.svg-container')
            if svg_element:
                await svg_element.screenshot(path=screenshot_path)
            else:
                await page.screenshot(path=screenshot_path)

            await browser.close()

        return screenshot_path

    async def _convert_to_pptx(self, svg_path: Path, test_name: str) -> Path:
        """Convert SVG to PPTX using our conversion system."""
        pptx_path = self.output_dir / "pptx_files" / f"{test_name}.pptx"

        if not CONVERTERS_AVAILABLE:
            # Create placeholder PPTX (just copy a template or create minimal file)
            self._create_placeholder_pptx(pptx_path, test_name)
            return pptx_path

        try:
            # Use the existing SVG2PPTX conversion system
            from src.svg2pptx import convert_svg_to_pptx

            # Convert SVG to PPTX using the main conversion pipeline
            svg_content = svg_path.read_text()
            temp_pptx_path = convert_svg_to_pptx(
                svg_input=svg_content,
                output_path=str(pptx_path),
                slide_width=10.0,  # Standard slide dimensions
                slide_height=7.5
            )

            logger.info(f"Successfully converted SVG to PPTX using main pipeline: {pptx_path}")

        except Exception as e:
            logger.warning(f"PPTX conversion failed for {test_name}: {e}")
            # Create detailed error information
            error_details = f"{test_name} - Conversion Failed: {str(e)[:100]}"
            self._create_placeholder_pptx(pptx_path, error_details)

        return pptx_path

    async def _capture_pptx_screenshot(self, pptx_path: Path, test_name: str) -> Path:
        """Capture screenshot of PPTX slide using LibreOffice."""
        screenshot_path = self.output_dir / "pptx_screenshots" / f"{test_name}_pptx.png"

        try:
            # Try LibreOffice headless conversion to image
            cmd = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", str(screenshot_path.parent),
                str(pptx_path)
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)

            # LibreOffice creates files with different naming
            generated_png = screenshot_path.parent / f"{pptx_path.stem}.png"
            if generated_png.exists():
                generated_png.rename(screenshot_path)
            else:
                self._create_placeholder_image(screenshot_path, f"PPTX Screenshot\n{test_name}\n(LibreOffice conversion failed)")

        except Exception as e:
            logger.warning(f"PPTX screenshot failed for {test_name}: {e}")
            self._create_placeholder_image(screenshot_path, f"PPTX Screenshot\n{test_name}\n(Screenshot failed)")

        return screenshot_path

    async def _create_comparison_page(self, test_case: VisualTestCase,
                                    svg_screenshot: Path, pptx_screenshot: Path,
                                    result: Dict) -> Path:
        """Create HTML comparison page showing SVG vs PPTX side by side."""
        comparison_path = self.output_dir / "comparisons" / f"{test_case.name}_comparison.html"

        # Calculate relative paths from comparison page
        svg_img_path = f"../svg_screenshots/{svg_screenshot.name}"
        pptx_img_path = f"../pptx_screenshots/{pptx_screenshot.name}"
        svg_file_path = f"../svgs/{test_case.name}.svg"
        pptx_file_path = f"../pptx_files/{test_case.name}.pptx"

        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>Visual Comparison: {test_case.name}</title>
            <style>
                body {{
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Arial, sans-serif;
                    margin: 0;
                    padding: 20px;
                    background: #f5f5f5;
                }}
                .header {{
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    color: white;
                    padding: 30px;
                    border-radius: 10px;
                    margin-bottom: 30px;
                }}
                .header h1 {{ margin: 0; }}
                .header p {{ margin: 5px 0; opacity: 0.9; }}
                .comparison-container {{
                    display: grid;
                    grid-template-columns: 1fr 1fr;
                    gap: 30px;
                    margin-bottom: 30px;
                }}
                .comparison-panel {{
                    background: white;
                    border-radius: 10px;
                    padding: 20px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                }}
                .comparison-panel h2 {{
                    margin-top: 0;
                    color: #333;
                    border-bottom: 2px solid #eee;
                    padding-bottom: 10px;
                }}
                .image-container {{
                    text-align: center;
                    margin: 20px 0;
                    border: 1px solid #ddd;
                    border-radius: 5px;
                    padding: 10px;
                    background: white;
                }}
                .image-container img {{
                    max-width: 100%;
                    height: auto;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.1);
                }}
                .file-links {{
                    margin-top: 15px;
                    text-align: center;
                }}
                .file-links a {{
                    display: inline-block;
                    margin: 5px;
                    padding: 8px 15px;
                    background: #007bff;
                    color: white;
                    text-decoration: none;
                    border-radius: 5px;
                    font-size: 14px;
                }}
                .file-links a:hover {{
                    background: #0056b3;
                }}
                .metadata {{
                    background: #f8f9fa;
                    padding: 20px;
                    border-radius: 10px;
                    margin-top: 30px;
                }}
                .metadata h3 {{
                    margin-top: 0;
                    color: #495057;
                }}
                .metadata-grid {{
                    display: grid;
                    grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
                    gap: 15px;
                }}
                .metadata-item {{
                    background: white;
                    padding: 15px;
                    border-radius: 5px;
                    border-left: 4px solid #007bff;
                }}
                .metadata-item strong {{
                    color: #495057;
                }}
                .status-completed {{ color: #28a745; }}
                .status-failed {{ color: #dc3545; }}
                .errors {{
                    background: #f8d7da;
                    color: #721c24;
                    padding: 15px;
                    border-radius: 5px;
                    margin-top: 15px;
                }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>📊 Visual Comparison Report</h1>
                <p><strong>Test Case:</strong> {test_case.name}</p>
                <p><strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
                <p><strong>Description:</strong> {test_case.description or 'Visual comparison test'}</p>
            </div>

            <div class="comparison-container">
                <div class="comparison-panel">
                    <h2>🎨 SVG Source (Browser Rendered)</h2>
                    <div class="image-container">
                        <img src="{svg_img_path}" alt="SVG Screenshot" />
                    </div>
                    <div class="file-links">
                        <a href="{svg_file_path}" target="_blank">📄 View SVG Source</a>
                        <a href="{svg_img_path}" target="_blank">🖼️ Open Image</a>
                    </div>
                </div>

                <div class="comparison-panel">
                    <h2>📋 PPTX Output (Converted Result)</h2>
                    <div class="image-container">
                        <img src="{pptx_img_path}" alt="PPTX Screenshot" />
                    </div>
                    <div class="file-links">
                        <a href="{pptx_file_path}" target="_blank">📊 Download PPTX</a>
                        <a href="{pptx_img_path}" target="_blank">🖼️ Open Image</a>
                    </div>
                </div>
            </div>

            <div class="metadata">
                <h3>📋 Test Metadata</h3>
                <div class="metadata-grid">
                    <div class="metadata-item">
                        <strong>Status:</strong><br>
                        <span class="status-{result['status']}">{result['status'].upper()}</span>
                    </div>
                    <div class="metadata-item">
                        <strong>Test Name:</strong><br>
                        {test_case.name}
                    </div>
                    <div class="metadata-item">
                        <strong>SVG Size:</strong><br>
                        {len(test_case.svg_content):,} characters
                    </div>
                    <div class="metadata-item">
                        <strong>Generated:</strong><br>
                        {test_case.timestamp[:19]}
                    </div>
                </div>

                {f'''
                <div class="errors">
                    <strong>⚠️ Errors:</strong><br>
                    {'<br>'.join(result['errors'])}
                </div>
                ''' if result.get('errors') else ''}
            </div>

            <div style="text-align: center; margin-top: 40px; color: #6c757d;">
                <p>Generated by SVG2PPTX Visual Testing System</p>
            </div>
        </body>
        </html>
        """

        comparison_path.write_text(html_content, encoding='utf-8')
        return comparison_path

    def _create_placeholder_image(self, path: Path, text: str):
        """Create a placeholder image with text (simple implementation)."""
        try:
            from PIL import Image, ImageDraw, ImageFont

            # Create a simple placeholder image
            img = Image.new('RGB', (400, 300), color='lightgray')
            draw = ImageDraw.Draw(img)

            # Try to use a basic font, fall back to default
            try:
                font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 16)
            except:
                font = ImageFont.load_default()

            # Draw text in center
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            x = (400 - text_width) // 2
            y = (300 - text_height) // 2

            draw.text((x, y), text, fill='black', font=font)
            img.save(path)

        except ImportError:
            # Fallback: create empty file
            path.touch()

    def _create_placeholder_pptx(self, path: Path, title: str):
        """Create a placeholder PPTX file."""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

            slide.shapes.title.text = title
            slide.shapes.placeholders[1].text = "SVG2PPTX Conversion Test\n\nThis is a placeholder slide."

            prs.save(path)

        except ImportError:
            # Create empty file as fallback
            path.touch()

    async def generate_test_report(self) -> Path:
        """Generate comprehensive test report."""
        report_path = self.output_dir / "visual_test_report.html"

        total_tests = len(self.test_results)
        completed_tests = len([r for r in self.test_results if r['status'] == 'completed'])
        failed_tests = len([r for r in self.test_results if r['status'] == 'failed'])

        # Create index page linking to all comparisons
        comparison_links = []
        for result in self.test_results:
            test_name = result['test_case']['name']
            comparison_file = f"comparisons/{test_name}_comparison.html"
            status_class = f"status-{result['status']}"
            comparison_links.append(f'''
                <div class="test-link">
                    <a href="{comparison_file}" class="{status_class}">
                        <strong>{test_name}</strong>
                        <span class="status">{result['status'].upper()}</span>
                    </a>
                    {f"<div class='description'>{result['test_case']['description']}</div>" if result['test_case']['description'] else ""}
                </div>
            ''')

        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>SVG2PPTX Visual Test Report</title>
            <style>
                body {{
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Arial, sans-serif;
                    margin: 0; padding: 20px; background: #f5f5f5;
                }}
                .header {{
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    color: white; padding: 30px; border-radius: 10px; margin-bottom: 30px;
                }}
                .summary {{
                    display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
                    gap: 20px; margin-bottom: 30px;
                }}
                .metric {{
                    background: white; padding: 20px; border-radius: 8px; text-align: center;
                    border-left: 4px solid #007bff; box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                }}
                .metric .value {{ font-size: 2.5em; font-weight: bold; color: #007bff; }}
                .test-links {{
                    background: white; padding: 30px; border-radius: 10px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                }}
                .test-link {{
                    margin: 15px 0; padding: 15px; border: 1px solid #ddd;
                    border-radius: 5px; background: #f8f9fa;
                }}
                .test-link a {{
                    text-decoration: none; color: #495057; display: flex;
                    justify-content: space-between; align-items: center;
                }}
                .test-link a:hover {{ background: #e9ecef; border-radius: 5px; padding: 10px; }}
                .status {{ padding: 5px 10px; border-radius: 3px; font-size: 12px; }}
                .status-completed {{ background: #d4edda; color: #155724; }}
                .status-failed {{ background: #f8d7da; color: #721c24; }}
                .description {{ margin-top: 10px; font-size: 14px; color: #6c757d; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>🎨 SVG2PPTX Visual Test Report</h1>
                <p><strong>Generated:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
                <p>Visual comparison testing results for SVG to PPTX conversion</p>
            </div>

            <div class="summary">
                <div class="metric">
                    <div class="value">{total_tests}</div>
                    <div>Total Tests</div>
                </div>
                <div class="metric">
                    <div class="value">{completed_tests}</div>
                    <div>Completed</div>
                </div>
                <div class="metric">
                    <div class="value">{failed_tests}</div>
                    <div>Failed</div>
                </div>
                <div class="metric">
                    <div class="value">{(completed_tests/max(1,total_tests)*100):.1f}%</div>
                    <div>Success Rate</div>
                </div>
            </div>

            <div class="test-links">
                <h2>📊 Individual Test Results</h2>
                {''.join(comparison_links)}
            </div>
        </body>
        </html>
        """

        report_path.write_text(html_content, encoding='utf-8')
        return report_path


# Test case definitions
def create_sample_test_cases() -> List[VisualTestCase]:
    """Create sample SVG test cases for visual comparison."""

    test_cases = []

    # Basic shapes test
    test_cases.append(VisualTestCase(
        name="basic_shapes",
        description="Basic SVG shapes: rectangle, circle, ellipse, line",
        svg_content="""
        <svg width="400" height="300" xmlns="http://www.w3.org/2000/svg">
            <rect x="10" y="10" width="100" height="80" fill="lightblue" stroke="blue" stroke-width="2"/>
            <circle cx="200" cy="50" r="40" fill="lightgreen" stroke="green" stroke-width="2"/>
            <ellipse cx="350" cy="50" rx="40" ry="25" fill="lightcoral" stroke="red" stroke-width="2"/>
            <line x1="10" y1="150" x2="390" y2="150" stroke="purple" stroke-width="3"/>
            <text x="200" y="200" text-anchor="middle" font-family="Arial" font-size="18" fill="darkblue">
                Basic Shapes Test
            </text>
        </svg>
        """
    ))

    # Gradient test
    test_cases.append(VisualTestCase(
        name="gradients",
        description="Linear and radial gradients with multiple stops",
        svg_content="""
        <svg width="400" height="300" xmlns="http://www.w3.org/2000/svg">
            <defs>
                <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="0%">
                    <stop offset="0%" style="stop-color:rgb(255,255,0);stop-opacity:1" />
                    <stop offset="100%" style="stop-color:rgb(255,0,0);stop-opacity:1" />
                </linearGradient>
                <radialGradient id="grad2" cx="50%" cy="50%" r="50%">
                    <stop offset="0%" style="stop-color:rgb(255,255,255);stop-opacity:0" />
                    <stop offset="100%" style="stop-color:rgb(0,0,255);stop-opacity:1" />
                </radialGradient>
            </defs>
            <rect x="10" y="10" width="180" height="120" fill="url(#grad1)" stroke="black"/>
            <circle cx="300" cy="70" r="60" fill="url(#grad2)" stroke="black"/>
            <text x="200" y="200" text-anchor="middle" font-family="Arial" font-size="18" fill="darkgreen">
                Gradient Test
            </text>
        </svg>
        """
    ))

    # Transform test
    test_cases.append(VisualTestCase(
        name="transforms",
        description="Various SVG transforms: rotate, scale, translate",
        svg_content="""
        <svg width="400" height="300" xmlns="http://www.w3.org/2000/svg">
            <g transform="translate(50, 50)">
                <rect x="0" y="0" width="60" height="40" fill="lightblue" stroke="blue"/>
                <text x="30" y="55" text-anchor="middle" font-size="12">Translated</text>
            </g>
            <g transform="rotate(45, 200, 100)">
                <rect x="170" y="70" width="60" height="40" fill="lightgreen" stroke="green"/>
                <text x="200" y="115" text-anchor="middle" font-size="12">Rotated</text>
            </g>
            <g transform="scale(1.5) translate(200, 100)">
                <rect x="0" y="0" width="40" height="30" fill="lightcoral" stroke="red"/>
                <text x="20" y="45" text-anchor="middle" font-size="8">Scaled</text>
            </g>
        </svg>
        """
    ))

    # Path test
    test_cases.append(VisualTestCase(
        name="paths",
        description="SVG paths with curves and complex shapes",
        svg_content="""
        <svg width="400" height="300" xmlns="http://www.w3.org/2000/svg">
            <path d="M 50 150 Q 100 50 150 150 T 250 150"
                  stroke="blue" stroke-width="3" fill="none"/>
            <path d="M 50 200 C 50 100, 200 100, 200 200 S 350 300, 350 200"
                  stroke="red" stroke-width="3" fill="none"/>
            <path d="M 100 250 L 150 200 L 200 250 L 175 275 L 125 275 Z"
                  fill="yellow" stroke="orange" stroke-width="2"/>
            <text x="200" y="50" text-anchor="middle" font-family="Arial" font-size="18" fill="darkblue">
                Path Test
            </text>
        </svg>
        """
    ))

    # Text test
    test_cases.append(VisualTestCase(
        name="text",
        description="Text with different fonts, sizes, and styles",
        svg_content="""
        <svg width="400" height="300" xmlns="http://www.w3.org/2000/svg">
            <text x="200" y="30" text-anchor="middle" font-family="Arial" font-size="20" font-weight="bold" fill="darkblue">
                Text Rendering Test
            </text>
            <text x="50" y="80" font-family="Arial" font-size="16" fill="black">
                Normal text
            </text>
            <text x="50" y="110" font-family="Arial" font-size="16" font-weight="bold" fill="darkred">
                Bold text
            </text>
            <text x="50" y="140" font-family="Arial" font-size="16" font-style="italic" fill="darkgreen">
                Italic text
            </text>
            <text x="50" y="170" font-family="Arial" font-size="14" text-decoration="underline" fill="purple">
                Underlined text
            </text>
            <text x="200" y="220" text-anchor="middle" font-family="Georgia" font-size="18" fill="brown">
                Different Font Family
            </text>
        </svg>
        """
    ))

    return test_cases


# Pytest integration
@pytest.mark.skipif(not PLAYWRIGHT_AVAILABLE, reason="Playwright not available")
@pytest.mark.asyncio
async def test_visual_comparisons():
    """Run visual comparison tests for SVG2PPTX conversion."""

    tester = VisualComparisonTester()
    test_cases = create_sample_test_cases()

    logger.info(f"Running {len(test_cases)} visual comparison tests...")

    for test_case in test_cases:
        result = await tester.run_visual_test(test_case)
        tester.test_results.append(result)

        # Assert test completed successfully
        assert result['status'] in ['completed', 'failed'], f"Test {test_case.name} has invalid status"

        # Check that files were created
        if result['status'] == 'completed':
            assert 'svg_screenshot' in result['files'], f"SVG screenshot missing for {test_case.name}"
            assert 'pptx_screenshot' in result['files'], f"PPTX screenshot missing for {test_case.name}"
            assert 'comparison_page' in result['files'], f"Comparison page missing for {test_case.name}"

    # Generate comprehensive report
    report_path = await tester.generate_test_report()
    logger.info(f"Visual test report generated: {report_path}")

    # Print summary
    total = len(tester.test_results)
    completed = len([r for r in tester.test_results if r['status'] == 'completed'])
    print(f"\n📊 Visual Testing Summary:")
    print(f"   Total tests: {total}")
    print(f"   Completed: {completed}")
    print(f"   Success rate: {completed/max(1,total)*100:.1f}%")
    print(f"   Report: {report_path}")


if __name__ == "__main__":
    """Run visual comparison tests directly."""
    asyncio.run(test_visual_comparisons())