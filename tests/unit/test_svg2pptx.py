#!/usr/bin/env python3
"""
Unit Tests for SVG2PPTX Main Module

Comprehensive tests for the main SVG to PPTX conversion functionality
including file processing, conversion pipeline, and error handling.
"""

import pytest
from unittest.mock import Mock, patch, MagicMock, mock_open
from pathlib import Path
import sys
from lxml import etree as ET
import tempfile
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent / "src"))

from src.svg2pptx import convert_svg_to_pptx, SVGToPowerPointConverter, main

class TestSVGToPowerPointConverter:
    """Test cases for SVGToPowerPointConverter class."""

    def setup_method(self):
        """Set up test fixtures before each test method."""
        self.converter = SVGToPowerPointConverter()

    # Initialization Tests
    def test_initialization_default(self):
        """Test SVGToPowerPointConverter initialization with default parameters."""
        converter = SVGToPowerPointConverter()
        assert converter is not None
        # Test default parameters if they exist

    def test_initialization_with_options(self):
        """Test SVGToPowerPointConverter initialization with custom options."""
        converter = SVGToPowerPointConverter(slide_width=12, slide_height=9)
        assert converter is not None

    # File Processing Tests
    def test_load_svg_from_string(self):
        """Test loading SVG from string input."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="blue"/>
        </svg>'''

        result = self.converter.load_svg(svg_content)
        assert result is not None
        assert result.tag.endswith('svg')

    def test_load_svg_from_file(self):
        """Test loading SVG from file path."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="red"/>
        </svg>'''

        with tempfile.NamedTemporaryFile(mode='w', suffix='.svg', delete=False) as f:
            f.write(svg_content)
            temp_path = f.name

        try:
            result = self.converter.load_svg(temp_path)
            assert result is not None
            assert result.tag.endswith('svg')
        finally:
            Path(temp_path).unlink()

    def test_load_svg_invalid_file(self):
        """Test loading SVG from non-existent file."""
        with pytest.raises(FileNotFoundError):
            self.converter.load_svg('nonexistent.svg')

    def test_load_svg_invalid_xml(self):
        """Test loading invalid XML content."""
        invalid_svg = '<svg><rect></svg>'  # Missing closing tag
        with pytest.raises(ET.XMLSyntaxError):
            self.converter.load_svg(invalid_svg)

    def test_load_svg_empty_content(self):
        """Test loading empty SVG content."""
        with pytest.raises((ValueError, ET.XMLSyntaxError)):
            self.converter.load_svg('')

    # Conversion Tests
    @patch('src.svg2pptx.ConverterRegistry')
    @patch('src.svg2pptx.ConversionContext')
    def test_convert_svg_basic(self, mock_context, mock_registry):
        """Test basic SVG conversion."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="green"/>
        </svg>'''

        # Mock the converter registry and context
        mock_converter = Mock()
        mock_converter.convert.return_value = '<p:sp>mock_shape</p:sp>'
        mock_registry.return_value.get_converter.return_value = mock_converter
        mock_context.return_value = Mock()

        result = self.converter.convert(svg_content)
        assert result is not None
        assert isinstance(result, str)

    @patch('src.svg2pptx.ConverterRegistry')
    def test_convert_svg_no_elements(self, mock_registry):
        """Test conversion of SVG with no convertible elements."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <!-- Empty SVG -->
        </svg>'''

        mock_registry.return_value.get_converter.return_value = None

        result = self.converter.convert(svg_content)
        assert result is not None

    @patch('src.svg2pptx.ConverterRegistry')
    def test_convert_svg_multiple_elements(self, mock_registry):
        """Test conversion of SVG with multiple elements."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="30" height="30" fill="red"/>
            <circle cx="70" cy="70" r="15" fill="blue"/>
            <text x="50" y="50" fill="black">Test</text>
        </svg>'''

        mock_converter = Mock()
        mock_converter.convert.return_value = '<p:sp>mock_shape</p:sp>'
        mock_registry.return_value.get_converter.return_value = mock_converter

        result = self.converter.convert(svg_content)
        assert result is not None
        # Should have called converter for each element
        assert mock_converter.convert.call_count >= 1

    def test_drawingml_shapes_round_trip(self, tmp_path):
        """DrawingML snippets should materialize as shapes in the final PPTX."""

        svg_path = tmp_path / "sample.svg"
        svg_path.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10"/>',
            encoding='utf-8'
        )

        pptx_path = tmp_path / "output.pptx"

        drawingml_snippet = (
            "<p:sp>"
            "<p:nvSpPr>"
            "<p:cNvPr id=\"2\" name=\"Rectangle 1\"/>"
            "<p:cNvSpPr/>"
            "<p:nvPr/>"
            "</p:nvSpPr>"
            "<p:spPr>"
            "<a:xfrm>"
            "<a:off x=\"914400\" y=\"914400\"/>"
            "<a:ext cx=\"1828800\" cy=\"914400\"/>"
            "</a:xfrm>"
            "<a:prstGeom prst=\"rect\">"
            "<a:avLst/>"
            "</a:prstGeom>"
            "<a:solidFill>"
            "<a:srgbClr val=\"FF0000\"/>"
            "</a:solidFill>"
            "</p:spPr>"
            "<p:style>"
            "<a:lnRef idx=\"1\"><a:schemeClr val=\"accent1\"/></a:lnRef>"
            "<a:fillRef idx=\"3\"><a:schemeClr val=\"accent1\"/></a:fillRef>"
            "<a:effectRef idx=\"2\"><a:schemeClr val=\"accent1\"/></a:effectRef>"
            "<a:fontRef idx=\"minor\"><a:schemeClr val=\"lt1\"/></a:fontRef>"
            "</p:style>"
            "<p:txBody>"
            "<a:bodyPr rtlCol=\"0\" anchor=\"ctr\"/>"
            "<a:lstStyle/>"
            "<a:p><a:pPr algn=\"ctr\"/></a:p>"
            "</p:txBody>"
            "</p:sp>"
        )

        with patch('src.svg2pptx.SVGToDrawingMLConverter.convert_file', return_value=drawingml_snippet):
            result_path = self.converter.convert_file(str(svg_path), str(pptx_path))

        presentation = Presentation(result_path)
        slide = presentation.slides[0]
        non_placeholder_shapes = [
            shape for shape in slide.shapes if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER
        ]

        assert non_placeholder_shapes, "Expected DrawingML shapes to be added to the slide"
        assert any(
            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in non_placeholder_shapes
        ), "Inserted DrawingML shape should remain an auto shape after reload"

    # Error Handling Tests
    def test_convert_svg_conversion_error(self):
        """Test handling of conversion errors."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="yellow"/>
        </svg>'''

        with patch('src.svg2pptx.ConverterRegistry') as mock_registry:
            mock_converter = Mock()
            mock_converter.convert.side_effect = Exception("Conversion failed")
            mock_registry.return_value.get_converter.return_value = mock_converter

            with pytest.raises(Exception):
                self.converter.convert(svg_content)

    # Output Generation Tests
    @patch('src.svg2pptx.zipfile.ZipFile')
    def test_generate_pptx_output(self, mock_zipfile):
        """Test PPTX file generation."""
        mock_zip = Mock()
        mock_zipfile.return_value.__enter__.return_value = mock_zip

        drawingml_content = '<p:sp>mock_shape</p:sp>'
        output_path = 'test_output.pptx'

        self.converter.generate_pptx(drawingml_content, output_path)

        # Should have created zip file and written content
        mock_zipfile.assert_called_once()
        mock_zip.writestr.assert_called()

    def test_generate_pptx_invalid_path(self):
        """Test PPTX generation with invalid output path."""
        drawingml_content = '<p:sp>mock_shape</p:sp>'
        invalid_path = '/invalid/path/output.pptx'

        with pytest.raises((OSError, FileNotFoundError)):
            self.converter.generate_pptx(drawingml_content, invalid_path)

    # Viewport and Sizing Tests
    def test_extract_svg_dimensions(self):
        """Test extracting dimensions from SVG."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="200" height="150">
            <rect x="0" y="0" width="100%" height="100%" fill="white"/>
        </svg>'''

        svg_element = ET.fromstring(svg_content)
        width, height = self.converter.extract_dimensions(svg_element)

        assert width == 200
        assert height == 150

    def test_extract_svg_dimensions_with_units(self):
        """Test extracting dimensions with units from SVG."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="200px" height="150pt">
            <rect x="0" y="0" width="100%" height="100%" fill="white"/>
        </svg>'''

        svg_element = ET.fromstring(svg_content)
        width, height = self.converter.extract_dimensions(svg_element)

        assert width > 0
        assert height > 0

    def test_extract_svg_dimensions_viewbox(self):
        """Test extracting dimensions from SVG viewBox."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 200">
            <rect x="0" y="0" width="100%" height="100%" fill="white"/>
        </svg>'''

        svg_element = ET.fromstring(svg_content)
        width, height = self.converter.extract_dimensions(svg_element)

        assert width == 300
        assert height == 200

    def test_extract_svg_dimensions_default(self):
        """Test extracting dimensions when none specified."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg">
            <rect x="0" y="0" width="100" height="100" fill="white"/>
        </svg>'''

        svg_element = ET.fromstring(svg_content)
        width, height = self.converter.extract_dimensions(svg_element)

        # Should return reasonable defaults
        assert width > 0
        assert height > 0


class TestConvertSvgToPptx:
    """Test cases for convert_svg_to_pptx utility function."""

    def test_convert_function_with_file_paths(self):
        """Test convert function with file paths."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="purple"/>
        </svg>'''

        with tempfile.NamedTemporaryFile(mode='w', suffix='.svg', delete=False) as svg_file:
            svg_file.write(svg_content)
            svg_path = svg_file.name

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as pptx_file:
            pptx_path = pptx_file.name

        try:
            with patch('src.svg2pptx.SVG2PPTX') as mock_class:
                mock_instance = Mock()
                mock_class.return_value = mock_instance
                mock_instance.convert.return_value = '<p:sp>mock</p:sp>'

                result = convert_svg_to_pptx(svg_path, pptx_path)

                # Should have called the converter
                mock_class.assert_called_once()
                mock_instance.convert.assert_called_once()
                assert result is not None
        finally:
            Path(svg_path).unlink()
            Path(pptx_path).unlink(missing_ok=True)

    def test_convert_function_with_string_content(self):
        """Test convert function with string content."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <circle cx="50" cy="50" r="40" fill="orange"/>
        </svg>'''

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as pptx_file:
            pptx_path = pptx_file.name

        try:
            with patch('src.svg2pptx.SVG2PPTX') as mock_class:
                mock_instance = Mock()
                mock_class.return_value = mock_instance
                mock_instance.convert.return_value = '<p:sp>mock</p:sp>'

                result = convert_svg_to_pptx(svg_content, pptx_path)

                mock_class.assert_called_once()
                mock_instance.convert.assert_called_once()
                assert result is not None
        finally:
            Path(pptx_path).unlink(missing_ok=True)

    def test_convert_function_with_options(self):
        """Test convert function with custom options."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="cyan"/>
        </svg>'''

        options = {
            'background_color': 'white',
            'preserve_aspect_ratio': True
        }

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as pptx_file:
            pptx_path = pptx_file.name

        try:
            with patch('src.svg2pptx.SVG2PPTX') as mock_class:
                mock_instance = Mock()
                mock_class.return_value = mock_instance
                mock_instance.convert.return_value = '<p:sp>mock</p:sp>'

                result = convert_svg_to_pptx(svg_content, pptx_path, **options)

                # Should pass options to constructor
                mock_class.assert_called_once_with(**options)
                assert result is not None
        finally:
            Path(pptx_path).unlink(missing_ok=True)


class TestMain:
    """Test cases for main function (CLI interface)."""

    def test_main_function_basic(self):
        """Test main function with basic arguments."""
        test_args = ['svg2pptx', 'input.svg', 'output.pptx']

        with patch('sys.argv', test_args):
            with patch('src.svg2pptx.convert_svg_to_pptx') as mock_convert:
                with patch('src.svg2pptx.argparse.ArgumentParser.parse_args') as mock_args:
                    mock_args.return_value = Mock(
                        input='input.svg',
                        output='output.pptx',
                        width=None,
                        height=None,
                        background=None
                    )

                    main()

                    mock_convert.assert_called_once_with(
                        'input.svg',
                        'output.pptx'
                    )

    def test_main_function_with_options(self):
        """Test main function with command line options."""
        with patch('src.svg2pptx.convert_svg_to_pptx') as mock_convert:
            with patch('src.svg2pptx.argparse.ArgumentParser.parse_args') as mock_args:
                mock_args.return_value = Mock(
                    input='input.svg',
                    output='output.pptx',
                    width=1920,
                    height=1080,
                    background='white'
                )

                main()

                mock_convert.assert_called_once_with(
                    'input.svg',
                    'output.pptx',
                    width=1920,
                    height=1080,
                    background_color='white'
                )

    def test_main_function_error_handling(self):
        """Test main function error handling."""
        with patch('src.svg2pptx.convert_svg_to_pptx') as mock_convert:
            mock_convert.side_effect = Exception("Conversion failed")

            with patch('src.svg2pptx.argparse.ArgumentParser.parse_args') as mock_args:
                mock_args.return_value = Mock(
                    input='input.svg',
                    output='output.pptx',
                    width=None,
                    height=None,
                    background=None
                )

                with pytest.raises(SystemExit):
                    main()


# Integration Tests
class TestSVG2PPTXIntegration:
    """Integration tests for complete SVG to PPTX conversion workflow."""

    def test_complete_conversion_workflow(self):
        """Test complete conversion from SVG string to PPTX file."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="200" height="200">
            <rect x="50" y="50" width="100" height="100" fill="#FF0000"/>
            <circle cx="100" cy="100" r="30" fill="#00FF00"/>
            <text x="100" y="180" text-anchor="middle" fill="#0000FF">Test SVG</text>
        </svg>'''

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as pptx_file:
            pptx_path = pptx_file.name

        try:
            # This would test the actual conversion if modules are available
            # For now, we mock the heavy dependencies
            with patch('src.svg2pptx.ConverterRegistry') as mock_registry:
                mock_converter = Mock()
                mock_converter.convert.return_value = '<p:sp>converted_shape</p:sp>'
                mock_registry.return_value.get_converter.return_value = mock_converter

                converter = SVG2PPTX()
                result = converter.convert(svg_content)

                assert result is not None
                assert isinstance(result, str)
        finally:
            Path(pptx_path).unlink(missing_ok=True)

    def test_complex_svg_conversion(self):
        """Test conversion of complex SVG with nested groups."""
        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="300" height="300">
            <g transform="translate(50, 50)">
                <rect x="0" y="0" width="200" height="100" fill="#FFFF00" stroke="#000000"/>
                <g transform="rotate(45, 100, 50)">
                    <circle cx="100" cy="50" r="20" fill="#FF00FF"/>
                </g>
            </g>
            <path d="M10,10 L50,50 L10,90 Z" fill="#00FFFF"/>
        </svg>'''

        # Mock complex conversion scenario
        with patch('src.svg2pptx.ConverterRegistry') as mock_registry:
            mock_converter = Mock()
            mock_converter.convert.return_value = '<p:sp>complex_shape</p:sp>'
            mock_registry.return_value.get_converter.return_value = mock_converter

            converter = SVG2PPTX()
            result = converter.convert(svg_content)

            assert result is not None
            # Should handle nested groups and transformations
            assert mock_converter.convert.called


# Performance Tests
class TestSVG2PPTXPerformance:
    """Performance tests for SVG2PPTX conversion."""

    def setup_method(self):
        """Set up test fixtures before each test method."""
        self.converter = SVG2PPTX()

    @pytest.mark.performance
    def test_small_svg_performance(self):
        """Test performance with small SVG files."""
        import time

        svg_content = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="red"/>
        </svg>'''

        iterations = 100

        with patch('src.svg2pptx.ConverterRegistry') as mock_registry:
            mock_converter = Mock()
            mock_converter.convert.return_value = '<p:sp>shape</p:sp>'
            mock_registry.return_value.get_converter.return_value = mock_converter

            start_time = time.time()
            for _ in range(iterations):
                self.converter.convert(svg_content)
            end_time = time.time()

            duration = end_time - start_time
            assert duration < 5.0  # Should complete 100 conversions in less than 5 seconds

    @pytest.mark.performance
    def test_large_svg_performance(self):
        """Test performance with larger SVG files."""
        import time

        # Generate larger SVG with multiple elements
        elements = []
        for i in range(50):
            elements.append(f'<rect x="{i*10}" y="{i*10}" width="50" height="50" fill="#{i:02x}{i:02x}{i:02x}"/>')

        svg_content = f'''<svg xmlns="http://www.w3.org/2000/svg" width="1000" height="1000">
            {''.join(elements)}
        </svg>'''

        with patch('src.svg2pptx.ConverterRegistry') as mock_registry:
            mock_converter = Mock()
            mock_converter.convert.return_value = '<p:sp>shape</p:sp>'
            mock_registry.return_value.get_converter.return_value = mock_converter

            start_time = time.time()
            result = self.converter.convert(svg_content)
            end_time = time.time()

            duration = end_time - start_time
            assert duration < 10.0  # Should complete large conversion in less than 10 seconds
            assert result is not None


if __name__ == "__main__":
    pytest.main([__file__, "-v"])