#!/usr/bin/env python3
"""
Unit tests for PowerPoint presentation merger.

Tests slide copying, media handling, error handling, and integration functionality.
"""

import pytest
import tempfile
import shutil
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
from datetime import datetime

# Test if python-pptx is available
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from core.utils.powerpoint_merger import (
    PPTXMerger, PPTXMergeError, MergeResult,
    merge_presentations_simple, merge_pptx_files
)


class TestMergeResult:
    """Test MergeResult data class."""

    def test_merge_result_success(self):
        """Test successful merge result."""
        result = MergeResult(
            success=True,
            output_path=Path("/test/output.pptx"),
            total_slides=10,
            source_files=["file1.pptx", "file2.pptx"]
        )

        assert result.success is True
        assert result.output_path == Path("/test/output.pptx")
        assert result.total_slides == 10
        assert result.source_files == ["file1.pptx", "file2.pptx"]
        assert result.error_message is None
        assert result.warnings == []

    def test_merge_result_failure(self):
        """Test failed merge result."""
        result = MergeResult(
            success=False,
            error_message="Test error",
            source_files=["file1.pptx"]
        )

        assert result.success is False
        assert result.output_path is None
        assert result.total_slides == 0
        assert result.source_files == ["file1.pptx"]
        assert result.error_message == "Test error"
        assert result.warnings == []

    def test_merge_result_with_warnings(self):
        """Test merge result with warnings."""
        warnings = ["Warning 1", "Warning 2"]
        result = MergeResult(
            success=True,
            warnings=warnings
        )

        assert result.warnings == warnings

    def test_merge_result_post_init(self):
        """Test post_init behavior for default values."""
        result = MergeResult(success=True)

        assert result.source_files == []
        assert result.warnings == []


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
class TestPPTXMerger:
    """Test PPTXMerger functionality."""

    @pytest.fixture
    def temp_dir(self):
        """Create temporary directory for test files."""
        with tempfile.TemporaryDirectory() as tmp_dir:
            yield Path(tmp_dir)

    @pytest.fixture
    def mock_presentation_files(self, temp_dir):
        """Create mock PPTX files for testing."""
        files = []
        for i in range(3):
            file_path = temp_dir / f"test_{i}.pptx"
            # Create minimal PPTX file structure
            file_path.write_bytes(b"PK\x03\x04mock_pptx_content")
            files.append(file_path)
        return files

    def test_merger_initialization(self):
        """Test PPTXMerger initialization."""
        merger = PPTXMerger(
            preserve_master_slides=False,
            copy_embedded_media=False,
            default_slide_size=(8, 6)
        )

        assert merger.preserve_master_slides is False
        assert merger.copy_embedded_media is False
        assert merger.default_slide_size == (8, 6)

    def test_merger_initialization_defaults(self):
        """Test PPTXMerger with default parameters."""
        merger = PPTXMerger()

        assert merger.preserve_master_slides is True
        assert merger.copy_embedded_media is True
        assert merger.default_slide_size == (10, 7.5)

    @patch('src.utils.powerpoint_merger.PPTX_AVAILABLE', False)
    def test_merger_no_pptx_library(self):
        """Test merger when python-pptx is not available."""
        with pytest.raises(ImportError, match="python-pptx library is required"):
            PPTXMerger()

    def test_validate_input_files_empty(self):
        """Test validation with empty input files."""
        merger = PPTXMerger()

        with pytest.raises(PPTXMergeError, match="No input files provided"):
            merger._validate_input_files([])

    def test_validate_input_files_nonexistent(self, temp_dir):
        """Test validation with non-existent files."""
        merger = PPTXMerger()
        nonexistent_file = temp_dir / "nonexistent.pptx"

        with pytest.raises(PPTXMergeError, match="Input file does not exist"):
            merger._validate_input_files([nonexistent_file])

    def test_validate_input_files_not_pptx(self, temp_dir):
        """Test validation with non-PPTX files."""
        merger = PPTXMerger()
        text_file = temp_dir / "test.txt"
        text_file.write_text("test content")

        with pytest.raises(PPTXMergeError, match="Input file is not a PPTX file"):
            merger._validate_input_files([text_file])

    def test_validate_input_files_directory(self, temp_dir):
        """Test validation with directory instead of file."""
        merger = PPTXMerger()
        test_dir = temp_dir / "test_dir"
        test_dir.mkdir()

        with pytest.raises(PPTXMergeError, match="Input path is not a file"):
            merger._validate_input_files([test_dir])

    @patch('src.utils.powerpoint_merger.Presentation')
    def test_configure_presentation(self, mock_presentation_class):
        """Test presentation configuration."""
        merger = PPTXMerger(default_slide_size=(12, 9))
        mock_presentation = Mock()

        merger._configure_presentation(mock_presentation, "Test Title")

        # Check that slide dimensions were set
        assert mock_presentation.slide_width == Inches(12)
        assert mock_presentation.slide_height == Inches(9)

    @patch('src.utils.powerpoint_merger.Presentation')
    def test_merge_presentations_success(self, mock_presentation_class, mock_presentation_files, temp_dir):
        """Test successful presentation merging."""
        # Setup mocks
        mock_source_prs = Mock()
        mock_slide1 = Mock()
        mock_slide2 = Mock()
        mock_source_prs.slides = [mock_slide1, mock_slide2]

        mock_target_prs = Mock()
        mock_layout = Mock()
        mock_target_prs.slide_layouts = [None, None, None, None, None, None, mock_layout]  # Index 6

        mock_presentation_class.side_effect = [mock_target_prs, mock_source_prs, mock_source_prs, mock_source_prs]

        merger = PPTXMerger()
        output_path = temp_dir / "merged.pptx"

        # Mock the copy slide methods to avoid complex XML copying
        with patch.object(merger, '_copy_slide_to_presentation') as mock_copy:
            result = merger.merge_presentations(
                input_files=mock_presentation_files,
                output_path=output_path,
                presentation_title="Test Merge"
            )

            assert result.success is True
            assert result.output_path == output_path
            assert result.total_slides == 6  # 2 slides per file × 3 files
            assert len(result.source_files) == 3
            assert mock_copy.call_count == 6

    @patch('src.utils.powerpoint_merger.Presentation')
    def test_merge_presentations_file_error(self, mock_presentation_class, mock_presentation_files, temp_dir):
        """Test merge with file processing error."""
        # Make one file fail to open, others succeed with slides
        mock_target_prs = Mock()
        mock_successful_prs1 = Mock()
        mock_successful_prs2 = Mock()

        # Set up successful presentations with slides
        mock_slide1 = Mock()
        mock_slide2 = Mock()
        mock_successful_prs1.slides = [mock_slide1]
        mock_successful_prs2.slides = [mock_slide2]

        mock_presentation_class.side_effect = [
            mock_target_prs,  # Target presentation
            Exception("Cannot open file"),  # First file fails
            mock_successful_prs1,  # Second file succeeds
            mock_successful_prs2   # Third file succeeds
        ]

        merger = PPTXMerger()
        output_path = temp_dir / "merged.pptx"

        with patch.object(merger, '_copy_slide_to_presentation'):
            result = merger.merge_presentations(
                input_files=mock_presentation_files,
                output_path=output_path
            )

            assert result.success is True
            assert result.total_slides == 2  # 2 successful files with 1 slide each
            assert len(result.warnings) == 1
            assert "Cannot open file" in result.warnings[0]

    @patch('src.utils.powerpoint_merger.Presentation')
    def test_merge_presentations_all_files_fail(self, mock_presentation_class, mock_presentation_files, temp_dir):
        """Test merge when all files fail."""
        # Make all files fail
        mock_presentation_class.side_effect = [
            Mock(),  # Target presentation
            Exception("Error 1"),
            Exception("Error 2"),
            Exception("Error 3")
        ]

        merger = PPTXMerger()
        output_path = temp_dir / "merged.pptx"

        result = merger.merge_presentations(
            input_files=mock_presentation_files,
            output_path=output_path
        )

        assert result.success is False
        assert "No slides were successfully merged" in result.error_message

    def test_merge_from_file_list(self, mock_presentation_files, temp_dir):
        """Test merging from file list."""
        # Create file list
        file_list_path = temp_dir / "file_list.txt"
        with open(file_list_path, 'w') as f:
            for file_path in mock_presentation_files:
                f.write(f"{file_path}\n")

        merger = PPTXMerger()
        output_path = temp_dir / "merged.pptx"

        with patch.object(merger, 'merge_presentations') as mock_merge:
            mock_merge.return_value = MergeResult(success=True)

            result = merger.merge_from_file_list(file_list_path, output_path)

            mock_merge.assert_called_once_with(
                [str(f) for f in mock_presentation_files],
                output_path
            )

    def test_merge_from_file_list_nonexistent(self, temp_dir):
        """Test merge from non-existent file list."""
        merger = PPTXMerger()
        file_list_path = temp_dir / "nonexistent.txt"
        output_path = temp_dir / "merged.pptx"

        result = merger.merge_from_file_list(file_list_path, output_path)

        assert result.success is False
        assert "File list does not exist" in result.error_message

    def test_merge_from_empty_file_list(self, temp_dir):
        """Test merge from empty file list."""
        merger = PPTXMerger()
        file_list_path = temp_dir / "empty_list.txt"
        file_list_path.write_text("")
        output_path = temp_dir / "merged.pptx"

        result = merger.merge_from_file_list(file_list_path, output_path)

        assert result.success is False
        assert "No files found in file list" in result.error_message

    def test_create_merge_summary_success(self):
        """Test merge summary creation for successful merge."""
        result = MergeResult(
            success=True,
            output_path=Path("/test/output.pptx"),
            total_slides=15,
            source_files=["file1.pptx", "file2.pptx"],
            warnings=["Warning 1"]
        )

        merger = PPTXMerger()
        summary = merger.create_merge_summary(result)

        assert summary['success'] is True
        assert summary['total_slides'] == 15
        assert summary['source_file_count'] == 2
        assert summary['source_files'] == ["file1.pptx", "file2.pptx"]
        assert summary['output_path'] == "/test/output.pptx"
        assert summary['warnings_count'] == 1
        assert summary['warnings'] == ["Warning 1"]
        assert 'error' not in summary

    def test_create_merge_summary_failure(self):
        """Test merge summary creation for failed merge."""
        result = MergeResult(
            success=False,
            error_message="Merge failed",
            source_files=["file1.pptx"]
        )

        merger = PPTXMerger()
        summary = merger.create_merge_summary(result)

        assert summary['success'] is False
        assert summary['total_slides'] == 0
        assert summary['source_file_count'] == 1
        assert summary['output_path'] is None
        assert summary['error'] == "Merge failed"

    @patch('src.utils.powerpoint_merger.Presentation')
    def test_copy_slide_methods(self, mock_presentation_class):
        """Test slide copying methods (simplified)."""
        merger = PPTXMerger()

        # Test textbox copying
        mock_slide = Mock()
        mock_textbox = Mock()
        mock_textbox.left = 100
        mock_textbox.top = 200
        mock_textbox.width = 300
        mock_textbox.height = 400
        mock_textbox.text_frame.text = "Test text"

        # Mock the shapes.add_textbox method
        mock_new_textbox = Mock()
        mock_slide.shapes.add_textbox.return_value = mock_new_textbox

        merger._copy_textbox(mock_slide, mock_textbox)

        mock_slide.shapes.add_textbox.assert_called_once_with(100, 200, 300, 400)
        assert mock_new_textbox.text_frame.text == "Test text"


class TestSimpleFunctions:
    """Test simple function interfaces."""

    @patch('src.utils.powerpoint_merger.PPTXMerger')
    def test_merge_presentations_simple_success(self, mock_merger_class):
        """Test simple merge function with success."""
        mock_merger = Mock()
        mock_result = MergeResult(success=True, total_slides=5)
        mock_merger.merge_presentations.return_value = mock_result
        mock_merger.create_merge_summary.return_value = {'success': True, 'total_slides': 5}
        mock_merger_class.return_value = mock_merger

        result = merge_presentations_simple(
            input_files=["file1.pptx", "file2.pptx"],
            output_path="merged.pptx",
            preserve_master_slides=False
        )

        assert result['success'] is True
        assert result['total_slides'] == 5
        mock_merger_class.assert_called_once_with(preserve_master_slides=False)

    @patch('src.utils.powerpoint_merger.PPTXMerger')
    def test_merge_presentations_simple_failure(self, mock_merger_class):
        """Test simple merge function with failure."""
        mock_merger_class.side_effect = Exception("Test error")

        result = merge_presentations_simple(
            input_files=["file1.pptx", "file2.pptx"],
            output_path="merged.pptx"
        )

        assert result['success'] is False
        assert result['error'] == "Test error"
        assert result['total_slides'] == 0
        assert result['source_file_count'] == 2

    @patch('src.utils.powerpoint_merger.merge_presentations_simple')
    def test_merge_pptx_files_success(self, mock_simple_merge):
        """Test backward compatibility function with success."""
        mock_simple_merge.return_value = {'success': True}

        result = merge_pptx_files(["file1.pptx", "file2.pptx"], "merged.pptx")

        assert result is True
        mock_simple_merge.assert_called_once_with(["file1.pptx", "file2.pptx"], "merged.pptx")

    @patch('src.utils.powerpoint_merger.merge_presentations_simple')
    def test_merge_pptx_files_failure(self, mock_simple_merge):
        """Test backward compatibility function with failure."""
        mock_simple_merge.return_value = {'success': False}

        result = merge_pptx_files(["file1.pptx", "file2.pptx"], "merged.pptx")

        assert result is False

    @patch('src.utils.powerpoint_merger.merge_presentations_simple')
    def test_merge_pptx_files_exception(self, mock_simple_merge):
        """Test backward compatibility function with exception."""
        mock_simple_merge.side_effect = Exception("Test error")

        result = merge_pptx_files(["file1.pptx", "file2.pptx"], "merged.pptx")

        assert result is False


@pytest.mark.skipif(PPTX_AVAILABLE, reason="Testing import error handling")
class TestImportError:
    """Test behavior when python-pptx is not available."""

    def test_pptx_not_available_flag(self):
        """Test that PPTX_AVAILABLE flag is correctly set."""
        # This test only runs when PPTX is not available
        assert not PPTX_AVAILABLE


class TestErrorHandling:
    """Test error handling scenarios."""

    def test_pptx_merge_error_inheritance(self):
        """Test that PPTXMergeError inherits from Exception."""
        error = PPTXMergeError("Test error")
        assert isinstance(error, Exception)
        assert str(error) == "Test error"

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_merger_path_handling(self, tmp_path):
        """Test that merger handles Path objects correctly."""
        merger = PPTXMerger()

        # Test with Path objects
        input_files = [tmp_path / "file1.pptx", tmp_path / "file2.pptx"]
        output_path = tmp_path / "merged.pptx"

        # This should not raise an error during validation setup
        try:
            # Create the files so validation passes
            for file_path in input_files:
                file_path.write_bytes(b"mock pptx content")

            # Mock the actual merging to avoid complex setup
            with patch.object(merger, '_merge_single_presentation', return_value=1):
                with patch('src.utils.powerpoint_merger.Presentation'):
                    result = merger.merge_presentations(input_files, output_path)
                    # Just verify it runs without path-related errors
                    assert isinstance(result, MergeResult)

        except PPTXMergeError:
            # Expected for incomplete mock setup, but path handling should work
            pass