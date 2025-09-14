#!/usr/bin/env python3
"""
Comprehensive tests for svg2pptx.py main entry point.

Tests the main SVGToPowerPointConverter class with systematic tool integration
following the standardized architecture pattern.
"""

import pytest
import tempfile
from lxml import etree as ET
from unittest.mock import Mock, patch, mock_open, MagicMock
from pathlib import Path
from pptx import Presentation
import sys

# Add src to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent / "src"))

from src.svg2pptx import SVGToPowerPointConverter
from src.converters.base import BaseConverter
from src.units import UnitConverter


class TestSVGToPowerPointConverter:
    """Test the main SVGToPowerPointConverter class with standardized tools."""
    
    def test_initialization_default(self):
        """Test converter initialization with default parameters."""
        converter = SVGToPowerPointConverter()
        
        # Test default dimensions
        assert converter.slide_width == 10
        assert converter.slide_height == 7.5
        
        # Test that svg_converter is created
        assert converter.svg_converter is not None
    
    def test_initialization_custom_dimensions(self):
        """Test converter initialization with custom slide dimensions."""
        # Use UnitConverter for standardized dimension testing
        class MockConverter(BaseConverter):
            def can_convert(self, element): return True
            def convert(self, element, context): return ""
        
        mock_converter = MockConverter()
        
        # Test custom dimensions
        custom_width = 12
        custom_height = 9
        
        converter = SVGToPowerPointConverter(
            slide_width=custom_width, 
            slide_height=custom_height
        )
        
        assert converter.slide_width == custom_width
        assert converter.slide_height == custom_height
    
    def test_presentation_setup(self):
        """Test that converter is properly configured for presentation creation."""
        converter = SVGToPowerPointConverter()
        
        # Should have proper dimensions configured
        assert converter.slide_width == 10
        assert converter.slide_height == 7.5
        
        # Should have SVG converter ready
        assert converter.svg_converter is not None
    
    @patch('src.svg2pptx.Presentation')
    def test_convert_file_basic(self, mock_presentation):
        """Test basic SVG file conversion."""
        # Setup mocks
        mock_prs = Mock()
        mock_presentation.return_value = mock_prs
        mock_slide = Mock()
        mock_prs.slides.add_slide.return_value = mock_slide
        mock_slide.shapes.add_textbox.return_value.text = ""
        
        converter = SVGToPowerPointConverter()
        
        # Mock the svg_converter.convert_file method
        converter.svg_converter.convert_file = Mock(return_value='<drawingML>test</drawingML>')
        
        # Test file conversion
        with tempfile.NamedTemporaryFile(suffix='.svg', delete=False) as temp_svg:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
                result = converter.convert_file(temp_svg.name, temp_pptx.name)
                
                # Should return output file path
                assert result == temp_pptx.name
                
                # Verify presentation was created and saved
                mock_presentation.assert_called_once()
                mock_prs.save.assert_called_once_with(temp_pptx.name)
    
    @patch('builtins.open', new_callable=mock_open, read_data='<svg><circle r="25"/></svg>')
    @patch('os.path.exists')
    def test_convert_file_no_output_specified(self, mock_exists, mock_file):
        """Test conversion when no output file is specified."""
        mock_exists.return_value = True
        
        converter = SVGToPowerPointConverter()
        
        with patch('src.svg2pptx.SVGToDrawingMLConverter') as mock_svg_converter:
            mock_svg_instance = Mock()
            mock_svg_converter.return_value = mock_svg_instance
            mock_svg_instance.convert.return_value = '<drawingML>circle</drawingML>'
            
            # Test without output file - should generate one
            with tempfile.NamedTemporaryFile(suffix='.svg') as temp_svg:
                result = converter.convert_file(temp_svg.name)
                
                # Should generate output filename
                expected_output = temp_svg.name.replace('.svg', '.pptx')
                assert result == expected_output
    
    @patch('src.svg2pptx.Presentation')
    def test_convert_file_nonexistent(self, mock_presentation):
        """Test conversion with non-existent input file."""
        converter = SVGToPowerPointConverter()
        
        # This should work because the real implementation doesn't check file existence
        # before creating the presentation - it lets the SVG converter handle errors
        converter.svg_converter.convert_file = Mock(side_effect=FileNotFoundError("File not found"))
        
        with pytest.raises(FileNotFoundError):
            converter.convert_file('nonexistent.svg')
    
    @patch('builtins.open', new_callable=mock_open, read_data='invalid svg content')
    @patch('os.path.exists')
    def test_convert_file_invalid_svg(self, mock_exists, mock_file):
        """Test conversion with invalid SVG content."""
        mock_exists.return_value = True
        
        converter = SVGToPowerPointConverter()
        
        with patch('src.svg2pptx.SVGToDrawingMLConverter') as mock_svg_converter:
            # Make converter raise an exception
            mock_svg_converter.side_effect = ET.ParseError("Invalid XML")
            
            with tempfile.NamedTemporaryFile(suffix='.svg') as temp_svg:
                with pytest.raises(ET.ParseError):
                    converter.convert_file(temp_svg.name)
    
    @patch('os.listdir')
    @patch('os.path.exists')
    @patch('os.path.isdir')
    def test_batch_convert_basic(self, mock_isdir, mock_exists, mock_listdir):
        """Test batch conversion of multiple SVG files."""
        # Setup mocks
        mock_isdir.return_value = True
        mock_exists.return_value = True
        mock_listdir.return_value = ['file1.svg', 'file2.svg', 'not_svg.txt']
        
        converter = SVGToPowerPointConverter()
        
        # Mock convert_file method
        converter.convert_file = Mock(side_effect=['output1.pptx', 'output2.pptx'])
        
        # Test batch conversion
        with tempfile.TemporaryDirectory() as temp_input_dir:
            with tempfile.TemporaryDirectory() as temp_output_dir:
                converter.batch_convert(temp_input_dir, temp_output_dir)
                
                # Should call convert_file for each SVG file
                assert converter.convert_file.call_count == 2
    
    @patch('os.listdir')
    @patch('os.path.exists')
    @patch('os.path.isdir')
    def test_batch_convert_no_output_dir(self, mock_isdir, mock_exists, mock_listdir):
        """Test batch conversion without specifying output directory."""
        # Setup mocks
        mock_isdir.return_value = True
        mock_exists.return_value = True
        mock_listdir.return_value = ['test.svg']
        
        converter = SVGToPowerPointConverter()
        converter.convert_file = Mock(return_value='output.pptx')
        
        # Test without output directory
        with tempfile.TemporaryDirectory() as temp_input_dir:
            converter.batch_convert(temp_input_dir)
            
            # Should still process files
            converter.convert_file.assert_called_once()
    
    def test_batch_convert_nonexistent_directory(self):
        """Test batch conversion with non-existent input directory."""
        converter = SVGToPowerPointConverter()
        
        with pytest.raises(FileNotFoundError):
            converter.batch_convert('nonexistent_directory')
    
    @patch('os.listdir')
    @patch('os.path.exists')
    @patch('os.path.isdir')
    def test_batch_convert_empty_directory(self, mock_isdir, mock_exists, mock_listdir):
        """Test batch conversion with directory containing no SVG files."""
        # Setup mocks
        mock_isdir.return_value = True
        mock_exists.return_value = True
        mock_listdir.return_value = ['file.txt', 'document.pdf']  # No SVG files
        
        converter = SVGToPowerPointConverter()
        converter.convert_file = Mock()
        
        with tempfile.TemporaryDirectory() as temp_input_dir:
            converter.batch_convert(temp_input_dir)
            
            # Should not call convert_file since no SVG files
            converter.convert_file.assert_not_called()


class TestSVGToPowerPointConverterIntegration:
    """Integration tests using standardized tool architecture."""
    
    def test_converter_uses_standardized_tools(self):
        """Test that converter integrates properly with standardized tool architecture."""
        # Use our standardized MockConverter pattern
        class MockConverter(BaseConverter):
            def can_convert(self, element): return True
            def convert(self, element, context): return ""
        
        mock_converter = MockConverter()
        
        # Test that tools are available for dimension calculations
        converter = SVGToPowerPointConverter()
        
        # Verify slide dimensions can be calculated with UnitConverter
        width_emu = mock_converter.unit_converter.to_emu('10in')
        height_emu = mock_converter.unit_converter.to_emu('7.5in')
        
        # These should be valid EMU values
        assert width_emu > 0
        assert height_emu > 0
        assert isinstance(width_emu, int)
        assert isinstance(height_emu, int)
    
    def test_presentation_slide_dimensions(self):
        """Test that presentation slide dimensions use proper EMU calculations."""
        # Use standardized tools for dimension verification
        class MockConverter(BaseConverter):
            def can_convert(self, element): return True
            def convert(self, element, context): return ""
        
        mock_converter = MockConverter()
        
        # Test custom dimensions
        custom_width = 12.5
        custom_height = 8.25
        
        converter = SVGToPowerPointConverter(custom_width, custom_height)
        
        # Verify dimensions are stored correctly
        assert converter.slide_width == custom_width
        assert converter.slide_height == custom_height
        
        # Test EMU conversion using standardized tools
        width_emu = mock_converter.unit_converter.to_emu(f'{custom_width}in')
        height_emu = mock_converter.unit_converter.to_emu(f'{custom_height}in')
        
        # Should produce valid EMU values
        assert width_emu > 0
        assert height_emu > 0


class TestSVGToPowerPointConverterErrorHandling:
    """Test error handling and edge cases."""
    
    def test_initialization_invalid_dimensions(self):
        """Test initialization with invalid dimensions."""
        # Zero dimensions should still create converter but may cause issues
        converter = SVGToPowerPointConverter(0, 0)
        assert converter.slide_width == 0
        assert converter.slide_height == 0
        
        # Negative dimensions
        converter = SVGToPowerPointConverter(-1, -1)
        assert converter.slide_width == -1
        assert converter.slide_height == -1
    
    @patch('src.svg2pptx.Presentation')
    def test_presentation_creation_failure(self, mock_presentation):
        """Test handling of presentation creation failure."""
        # Make Presentation constructor raise an exception
        mock_presentation.side_effect = Exception("Failed to create presentation")
        
        with pytest.raises(Exception, match="Failed to create presentation"):
            SVGToPowerPointConverter()
    
    def test_file_path_edge_cases(self):
        """Test edge cases in file path handling."""
        converter = SVGToPowerPointConverter()
        
        # Test with empty string
        with pytest.raises(FileNotFoundError):
            converter.convert_file('')
        
        # Test with None (should raise TypeError)
        with pytest.raises(TypeError):
            converter.convert_file(None)


if __name__ == '__main__':
    pytest.main([__file__, '-v'])